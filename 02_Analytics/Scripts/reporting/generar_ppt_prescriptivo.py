"""
Genera outputs/presentations/20260625_Analitica_Prescriptiva_Molienda_T8.pptx
15 slides para audiencia ejecutiva/operacional (no Data Science).
"""

import os, sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm

ROOT = Path(__file__).parent.parent

# ── Paleta corporativa ─────────────────────────────────────────────────────────
AZUL       = RGBColor(0x1F, 0x38, 0x64)   # azul Codelco oscuro
AZUL_MED   = RGBColor(0x1A, 0x5E, 0x99)   # azul medio
VERDE      = RGBColor(0x27, 0xAE, 0x60)   # verde semaforo
NARANJA    = RGBColor(0xE6, 0x7E, 0x22)   # naranja semaforo
ROJO       = RGBColor(0xC0, 0x39, 0x2B)   # rojo semaforo
AMARILLO   = RGBColor(0xF3, 0x9C, 0x12)   # amarillo semaforo
GRIS_CLARO = RGBColor(0xF2, 0xF2, 0xF2)
BLANCO     = RGBColor(0xFF, 0xFF, 0xFF)
NEGRO      = RGBColor(0x1A, 0x1A, 0x1A)

# ── Rutas figuras ──────────────────────────────────────────────────────────────
FIG_CAUSAL   = ROOT / "outputs/figures/11_Modelo_Causal"
FIG_EVENT    = ROOT / "outputs/figures/02_EventStudy_T8"
FIG_RATES    = ROOT / "outputs/figures/06_Rates"
FIG_AUTON    = ROOT / "outputs/figures/05_Autonomia"

def fig(rel):
    """Resuelve ruta de figura, retorna None si no existe."""
    p = ROOT / rel
    return str(p) if p.exists() else None

# ── Helpers pptx ──────────────────────────────────────────────────────────────
def add_slide(prs, layout_idx=6):
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)

def fill_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=NEGRO,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

def add_img(slide, path, l, t, w, h=None):
    if path and Path(path).exists():
        if h:
            slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
        else:
            slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w))

def add_header_bar(slide, title, subtitle=None):
    """Barra azul oscura en top con titulo blanco."""
    add_rect(slide, 0, 0, 13.33, 1.1, AZUL)
    add_text(slide, title, 0.25, 0.08, 12, 0.7, size=24, bold=True,
             color=BLANCO, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.25, 0.72, 12, 0.4, size=13,
                 color=RGBColor(0xB0,0xC4,0xDE), align=PP_ALIGN.LEFT)

def add_footer(slide, text="División El Teniente · Analítica Prescriptiva Molienda · Junio 2026"):
    add_rect(slide, 0, 7.12, 13.33, 0.38, AZUL)
    add_text(slide, text, 0.2, 7.14, 12, 0.3, size=9,
             color=RGBColor(0xB0,0xC4,0xDE))

def bullet_block(slide, items, l, t, w, h, title=None, title_color=AZUL,
                 item_size=12, title_size=14):
    """Bloque de viñetas con título opcional."""
    offset = 0
    if title:
        add_text(slide, title, l, t, w, 0.35, size=title_size, bold=True,
                 color=title_color)
        offset = 0.35
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t+offset),
                                     Inches(w), Inches(h-offset))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(item_size)
        p.font.color.rgb = NEGRO
        p.space_before = Pt(3)

def semaforo_box(slide, l, t, label, value, color, size_val=20, size_lab=11):
    add_rect(slide, l, t, 1.7, 0.9, color)
    add_text(slide, value, l, t+0.05, 1.7, 0.55, size=size_val, bold=True,
             color=BLANCO, align=PP_ALIGN.CENTER)
    add_text(slide, label, l, t+0.58, 1.7, 0.35, size=size_lab,
             color=BLANCO, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDES
# ─────────────────────────────────────────────────────────────────────────────

def slide01_portada(prs):
    """S1 — Portada corporativa."""
    sl = add_slide(prs)
    fill_bg(sl, AZUL)

    # Franja superior degradada simulada
    add_rect(sl, 0, 0, 13.33, 3.5, AZUL_MED)
    add_rect(sl, 0, 0, 13.33, 0.12, RGBColor(0xE8, 0xA0, 0x00))  # línea dorada

    add_text(sl, "ANALÍTICA PRESCRIPTIVA", 0.6, 0.5, 12, 0.9,
             size=32, bold=True, color=BLANCO, align=PP_ALIGN.LEFT)
    add_text(sl, "OPTIMIZACIÓN DE MOLIENDA FRENTE A VENTANAS TENIENTE 8",
             0.6, 1.25, 11.5, 1.2, size=20, bold=False,
             color=RGBColor(0xB0,0xC4,0xDE), align=PP_ALIGN.LEFT)

    add_rect(sl, 0, 3.6, 13.33, 0.04, AMARILLO)

    add_text(sl, "División El Teniente · Superintendencia SAG", 0.6, 3.8, 10, 0.5,
             size=14, color=RGBColor(0xCC,0xDD,0xEE), align=PP_ALIGN.LEFT)
    add_text(sl, "Junio 2026", 0.6, 4.25, 4, 0.4,
             size=13, color=RGBColor(0xAA,0xBB,0xCC), align=PP_ALIGN.LEFT)

    add_text(sl, "Alcance: SAG1 · SAG2 · PMC · UNITARIO\n70 eventos T8 analizados · 93,612 registros 5-min · Sistema RT 3 capas",
             0.6, 5.0, 11, 0.9, size=12,
             color=RGBColor(0x90,0xA8,0xC0), align=PP_ALIGN.LEFT)

    add_rect(sl, 0, 7.12, 13.33, 0.38, RGBColor(0x12, 0x26, 0x44))
    add_text(sl, "CONFIDENCIAL — Uso interno División El Teniente",
             0.3, 7.15, 10, 0.3, size=9,
             color=RGBColor(0x80,0x90,0xA0))


def slide02_resumen_ejecutivo(prs):
    """S2 — Resumen ejecutivo: 5 hallazgos + 5 acciones."""
    sl = add_slide(prs)
    fill_bg(sl, GRIS_CLARO)
    add_header_bar(sl, "Resumen Ejecutivo",
                   "Cinco hallazgos críticos · Cinco acciones inmediatas")

    # Panel izquierdo — hallazgos
    add_rect(sl, 0.2, 1.2, 6.3, 5.7, BLANCO)
    add_text(sl, "HALLAZGOS CRÍTICOS", 0.35, 1.3, 6.0, 0.4,
             size=13, bold=True, color=AZUL)

    hallazgos = [
        "① Correa 315 parada el 49% del tiempo → déficit crónico pila SAG1",
        "② Autonomía SAG1 media 1.7 h (P10 = 0.5 h) · SAG2 media 2.6 h (P10 = 0.2 h)",
        "③ Riesgo real de agotamiento inicia en pila ≈ 18% (no 70% como se asumía)",
        "④ SAG1 pierde hasta 18% TPH en las 4 h post-T8 vs. 9% SAG2",
        "⑤ Sistema RT prescribe rates óptimos en <1 ms · sin recalcular modelos",
    ]
    for i, h in enumerate(hallazgos):
        add_text(sl, h, 0.35, 1.75+i*0.9, 6.0, 0.85,
                 size=11.5, color=NEGRO)
        add_rect(sl, 0.25, 1.75+i*0.9+0.02, 0.06, 0.55,
                 ROJO if i < 3 else NARANJA)

    # Panel derecho — acciones
    add_rect(sl, 6.8, 1.2, 6.3, 5.7, AZUL)
    add_text(sl, "ACCIONES INMEDIATAS", 7.0, 1.3, 6.0, 0.4,
             size=13, bold=True, color=BLANCO)

    acciones = [
        "① Adoptar tabla de rates prescritos por régimen (disponible hoy)",
        "② Fijar alarma en pila SAG1 < 18% · SAG2 < 18% (no 70%)",
        "③ Activar protocolo PRE-T8: aumentar rates 4 h antes",
        "④ Investigar y corregir disponibilidad correa 315",
        "⑤ Integrar semáforo IRO en sala de control (Power BI)",
    ]
    for i, a in enumerate(acciones):
        add_text(sl, a, 7.0, 1.75+i*0.9, 5.9, 0.85,
                 size=11.5, color=BLANCO)
        add_rect(sl, 6.9, 1.75+i*0.9+0.02, 0.06, 0.55, VERDE)

    add_footer(sl)


def slide03_problema_operacional(prs):
    """S3 — El problema: qué pasa cuando entra T8."""
    sl = add_slide(prs)
    fill_bg(sl, GRIS_CLARO)
    add_header_bar(sl, "El Problema Operacional",
                   "¿Qué ocurre con la molienda cuando entra una Ventana Teniente 8?")

    # Figura event study SAG1
    img1 = fig("outputs/figures/02_EventStudy_T8/advanced_t8_historical/01_EventStudy_SAG1.png")
    img2 = fig("outputs/figures/02_EventStudy_T8/advanced_t8_historical/02_EventStudy_SAG2.png")

    add_img(sl, img1, 0.2, 1.25, 6.3, 3.8)
    add_img(sl, img2, 6.8, 1.25, 6.3, 3.8)

    add_text(sl, "SAG1 — caída inmediata y lenta recuperación",
             0.2, 5.1, 6.3, 0.4, size=11, bold=True, color=AZUL,
             align=PP_ALIGN.CENTER)
    add_text(sl, "SAG2 — efecto diferido, mayor resiliencia de pila",
             6.8, 5.1, 6.3, 0.4, size=11, bold=True, color=AZUL_MED,
             align=PP_ALIGN.CENTER)

    # KPIs rápidos
    add_rect(sl, 0.2, 5.55, 12.9, 1.35, AZUL)
    kpis = [
        ("70 eventos", "T8 analizados"),
        ("93.612", "registros 5 min"),
        ("−18%", "TPH SAG1 pico T8"),
        ("−9%", "TPH SAG2 pico T8"),
        ("4–6 h", "recuperación SAG1"),
        ("2–4 h", "recuperación SAG2"),
    ]
    for i, (v, l) in enumerate(kpis):
        x = 0.25 + i * 2.15
        add_text(sl, v, x, 5.6, 2.0, 0.55, size=17, bold=True,
                 color=AMARILLO, align=PP_ALIGN.CENTER)
        add_text(sl, l, x, 6.1, 2.0, 0.35, size=9,
                 color=BLANCO, align=PP_ALIGN.CENTER)

    add_footer(sl)


def slide04_cadena_causal(prs):
    """S4 — Cadena causal con figura."""
    sl = add_slide(prs)
    fill_bg(sl, GRIS_CLARO)
    add_header_bar(sl, "Cadena Causal Operacional",
                   "Mecanismo validado en datos: de la parada de correa al agotamiento de pila")

    img = fig("outputs/figures/11_Modelo_Causal/01_Cadena_Causal_T8.png")
    if img:
        add_img(sl, img, 0.3, 1.25, 12.7, 5.5)
    else:
        # Diagrama de texto como fallback
        cajas = [
            (0.3,  2.5, "CORREA 315\nparada 49% tiempo", ROJO),
            (2.6,  2.5, "PILA SAG1\ndéficit crónico", NARANJA),
            (4.9,  2.5, "AUTONOMÍA\n< 1.7 h media", NARANJA),
            (7.2,  2.5, "RÉGIMEN\nEMERGENCIA", ROJO),
            (9.5,  2.5, "TPH CAÍDA\n−18%  SAG1", ROJO),
        ]
        for x, y, txt, col in cajas:
            add_rect(sl, x, y, 2.1, 1.1, col)
            add_text(sl, txt, x+0.05, y+0.1, 2.0, 0.9,
                     size=11, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
        # flechas (texto)
        for x in [2.42, 4.72, 7.02, 9.32]:
            add_text(sl, "→", x, 2.85, 0.3, 0.4,
                     size=18, bold=True, color=AZUL, align=PP_ALIGN.CENTER)

    add_text(sl, "Restricción estructural: el 49% de indisponibilidad de correa 315 "
             "no puede resolverse con optimización de rates — requiere mejora de infraestructura.",
             0.3, 6.8, 12.7, 0.45, size=11, color=ROJO, bold=True,
             align=PP_ALIGN.CENTER)
    add_footer(sl)


def slide05_hallazgos_clave(prs):
    """S5 — Hallazgos numéricos clave con semaforos."""
    sl = add_slide(prs)
    fill_bg(sl, GRIS_CLARO)
    add_header_bar(sl, "Hallazgos Clave — Datos",
                   "Umbrales reales validados en 70 eventos históricos")

    # Fila 1 semáforos
    semaforos_top = [
        ("Correa 315\nparada", "49%", ROJO),
        ("Autonomía SAG1\nmedia", "1.7 h", NARANJA),
        ("Autonomía SAG2\nmedia", "2.6 h", AMARILLO),
        ("Caída TPH SAG1\npico T8", "−18%", ROJO),
        ("Caída TPH SAG2\npico T8", "−9%", NARANJA),
        ("Recup. SAG1\n(mediana)", "4.3 h", NARANJA),
        ("Recup. SAG2\n(mediana)", "2.9 h", AMARILLO),
    ]
    for i, (lbl, val, col) in enumerate(semaforos_top):
        x = 0.2 + i * 1.84
        semaforo_box(sl, x, 1.2, lbl, val, col, size_val=17, size_lab=9)

    # Umbrales reales
    add_rect(sl, 0.2, 2.3, 12.9, 0.04, AZUL)
    add_text(sl, "UMBRALES REALES (validados estadísticamente):",
             0.2, 2.4, 12.5, 0.35, size=13, bold=True, color=AZUL)

    umbrales = [
        ("SAG1", "Riesgo de agotamiento inicia en pila ≤ 18%  (operadores usaban 70% como umbral conservador)", ROJO),
        ("SAG2", "Riesgo de agotamiento inicia en pila ≤ 18%  (capacidad efectiva 32,009 ton)", NARANJA),
        ("Autonomía", "Alarma alta ≤ 2.5 h   ·   Alarma crítica ≤ 1.0 h", AMARILLO),
    ]
    for i, (activo, texto, col) in enumerate(umbrales):
        add_rect(sl, 0.2, 2.85+i*0.72, 1.3, 0.55, col)
        add_text(sl, activo, 0.2, 2.85+i*0.72, 1.3, 0.55,
                 size=12, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
        add_text(sl, texto, 1.6, 2.85+i*0.72, 11.5, 0.55,
                 size=11.5, color=NEGRO)

    # Distribución riesgo SAG1
    add_rect(sl, 0.2, 5.1, 12.9, 0.04, AZUL)
    add_text(sl, "DISTRIBUCIÓN RIESGO OPERACIONAL SAG1 (tiempo histórico):",
             0.2, 5.2, 12.5, 0.35, size=12, bold=True, color=AZUL)

    dist = [("VERDE\n6.7%", VERDE), ("AMARILLO\n58.7%", AMARILLO),
            ("NARANJA\n31.8%", NARANJA), ("ROJO\n2.8%", ROJO)]
    total_w = 12.9
    shares = [0.067, 0.587, 0.318, 0.028]
    x = 0.2
    for (lbl, col), sh in zip(dist, shares):
        w = total_w * sh
        add_rect(sl, x, 5.6, w, 0.8, col)
        if w > 0.8:
            add_text(sl, lbl, x, 5.62, w, 0.7,
                     size=9, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
        x += w

    add_footer(sl)


def slide06_impacto_t8(prs):
    """S6 — Impacto cuantificado por duración T8."""
    sl = add_slide(prs)
    fill_bg(sl, GRIS_CLARO)
    add_header_bar(sl, "Impacto Cuantificado por Duración T8",
                   "Dosis-respuesta validada: a mayor duración T8, mayor caída de TPH y pila")

    img1 = fig("outputs/figures/02_EventStudy_T8/event_study/05_EventStudy_2h.png")
    img2 = fig("outputs/figures/02_EventStudy_T8/event_study/06_EventStudy_4h.png")
    img3 = fig("outputs/figures/02_EventStudy_T8/Dosis_Respuesta_T8.png")

    add_img(sl, img1, 0.2, 1.25, 4.1, 3.3)
    add_img(sl, img2, 4.5, 1.25, 4.1, 3.3)
    add_img(sl, img3, 8.9, 1.25, 4.2, 3.3)

    add_text(sl, "T8 = 2 h", 0.2, 4.6, 4.1, 0.3, size=10, bold=True,
             color=AZUL_MED, align=PP_ALIGN.CENTER)
    add_text(sl, "T8 = 4 h", 4.5, 4.6, 4.1, 0.3, size=10, bold=True,
             color=AZUL_MED, align=PP_ALIGN.CENTER)
    add_text(sl, "Dosis-Respuesta (todos eventos)", 8.9, 4.6, 4.2, 0.3,
             size=10, bold=True, color=AZUL_MED, align=PP_ALIGN.CENTER)

    # tabla resumen
    add_rect(sl, 0.2, 5.0, 12.9, 0.4, AZUL)
    headers = ["Duración T8", "Registros", "Caída TPH SAG1", "Caída TPH SAG2",
               "P(agot.) SAG1", "Recuperación SAG1"]
    widths = [1.9, 1.9, 2.1, 2.1, 2.1, 2.8]
    x = 0.25
    for h, w in zip(headers, widths):
        add_text(sl, h, x, 5.05, w, 0.35, size=9, bold=True,
                 color=BLANCO, align=PP_ALIGN.CENTER)
        x += w

    rows = [
        ("2 h", "16,002", "−8%", "−4%", "12%", "2.1 h"),
        ("4 h", "41,085", "−18%", "−9%", "28%", "4.3 h"),
        ("8 h", "961",    "−22%", "−11%", "41%", "6.1 h"),
        ("12 h","6,054",  "−26%", "−14%", "55%", "8.2 h"),
    ]
    colors_row = [GRIS_CLARO, BLANCO, GRIS_CLARO, BLANCO]
    for ri, (row, bg) in enumerate(zip(rows, colors_row)):
        add_rect(sl, 0.2, 5.42+ri*0.42, 12.9, 0.42, bg)
        x = 0.25
        for val, w in zip(row, widths):
            c = ROJO if "−" in val and float(val.replace("−","").replace("%","").replace(" h","") or 0) > 10 else NEGRO
            add_text(sl, val, x, 5.45+ri*0.42, w, 0.38,
                     size=10, color=c, align=PP_ALIGN.CENTER)
            x += w

    add_footer(sl)


def slide07_autonomia_pilas(prs):
    """S7 — KPI Autonomía y estado de pilas."""
    sl = add_slide(prs)
    fill_bg(sl, GRIS_CLARO)
    add_header_bar(sl, "Autonomía de Pilas — KPI Crítico",
                   "Tiempo disponible antes del agotamiento: SAG1 y SAG2")

    img1 = fig("outputs/figures/11_Modelo_Causal/02_Pila_vs_Riesgo_SAG1.png")
    img2 = fig("outputs/figures/11_Modelo_Causal/04_Autonomia_vs_Riesgo.png")
    img3 = fig("outputs/figures/05_Autonomia/autonomia/Historico_Autonomia_SAG1_SAG2.png")

    add_img(sl, img1, 0.2, 1.25, 4.2, 3.4)
    add_img(sl, img2, 4.6, 1.25, 4.2, 3.4)
    add_img(sl, img3, 8.9, 1.25, 4.2, 3.4)

    add_text(sl, "Pila SAG1 vs Riesgo\n(umbral real = 18%)", 0.2, 4.7, 4.2, 0.4,
             size=10, bold=True, color=AZUL, align=PP_ALIGN.CENTER)
    add_text(sl, "Autonomía vs P(agotamiento)\n(alarma < 2.5 h)", 4.6, 4.7, 4.2, 0.4,
             size=10, bold=True, color=AZUL, align=PP_ALIGN.CENTER)
    add_text(sl, "Histórico Autonomía SAG1 vs SAG2", 8.9, 4.7, 4.2, 0.4,
             size=10, bold=True, color=AZUL, align=PP_ALIGN.CENTER)

    # Formula + parámetros
    add_rect(sl, 0.2, 5.2, 12.9, 1.6, AZUL)
    add_text(sl,
             "Fórmula:   Autonomía (h) = (Pila% − Pila_crítica%) ÷ Tasa_descarga (%/h)",
             0.4, 5.28, 12.5, 0.45, size=14, bold=True, color=BLANCO)
    params = [
        ("SAG1", "Pila crítica 15%", "Tasa descarga 23.76 %/h", "Cap. efectiva 4,575 ton"),
        ("SAG2", "Pila crítica 18.2%","Tasa descarga 6.18 %/h",  "Cap. efectiva 32,009 ton"),
    ]
    for ri, (act, p1, p2, p3) in enumerate(params):
        add_rect(sl, 0.25, 5.78+ri*0.48, 1.1, 0.4, AZUL_MED)
        add_text(sl, act, 0.25, 5.78+ri*0.48, 1.1, 0.4, size=12,
                 bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
        add_text(sl, f"{p1}   ·   {p2}   ·   {p3}",
                 1.5, 5.78+ri*0.48, 11.2, 0.4, size=11, color=BLANCO)

    add_footer(sl)


def slide08_validacion_reglas(prs):
    """S8 — Validación de las 15 reglas operacionales."""
    sl = add_slide(prs)
    fill_bg(sl, GRIS_CLARO)
    add_header_bar(sl, "Validación de Reglas Operacionales",
                   "15 reglas hipotéticas probadas contra 70 eventos históricos")

    img = fig("outputs/figures/11_Modelo_Causal/07_Validacion_Reglas.png")
    if img:
        add_img(sl, img, 0.2, 1.25, 7.5, 5.6)

    # Panel derecho — resumen reglas validadas
    add_rect(sl, 8.0, 1.25, 5.1, 5.6, AZUL)
    add_text(sl, "REGLAS VALIDADAS (extracto)", 8.1, 1.35, 4.9, 0.4,
             size=12, bold=True, color=BLANCO)

    reglas = [
        ("✓", "Pila SAG1 < 18% → P(agot) > 20%", VERDE),
        ("✓", "Autonomía < 2.5 h → alarma alta", VERDE),
        ("✓", "Autonomía < 1.0 h → alarma crítica", VERDE),
        ("✓", "correa_315=0 → régimen EMERGENCIA", VERDE),
        ("✓", "PRE-T8: aumentar rate SAG2 4 h antes", VERDE),
        ("✓", "DURANTE: reducir rate SAG1 si pila < 18%", VERDE),
        ("✓", "POST-T8: mantener rate elevado hasta +4 h", VERDE),
        ("~", "SAG1 recovery > 20% con solo rates", NARANJA),
        ("✗", "70% pila = umbral riesgo real (dato obsoleto)", ROJO),
    ]
    for i, (icono, texto, col) in enumerate(reglas):
        add_rect(sl, 8.05, 1.8+i*0.53, 0.45, 0.45, col)
        add_text(sl, icono, 8.05, 1.8+i*0.53, 0.45, 0.45,
                 size=13, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
        add_text(sl, texto, 8.6, 1.82+i*0.53, 4.4, 0.42,
                 size=9.5, color=BLANCO)

    add_footer(sl)


def slide09_matriz_riesgo(prs):
    """S9 — Matriz de riesgo pila × autonomía."""
    sl = add_slide(prs)
    fill_bg(sl, GRIS_CLARO)
    add_header_bar(sl, "Matriz de Riesgo Operacional",
                   "Estado combinado: nivel de pila × autonomía disponible")

    img = fig("outputs/figures/11_Modelo_Causal/08_Heatmap_Riesgo.png")
    img2 = fig("outputs/figures/11_Modelo_Causal/ejecutivo/E10_Matriz_Riesgo.png")

    if img:
        add_img(sl, img, 0.2, 1.25, 8.2, 5.1)
    if img2:
        add_img(sl, img2, 8.5, 1.25, 4.6, 5.1)

    # Leyenda semáforo
    legend = [
        ("VERDE  > 80", "Operación normal, sin restricción", VERDE),
        ("AMARILLO 60-80", "Monitoreo activo, precaución", AMARILLO),
        ("NARANJA 40-60", "Activar protocolo T8 preventivo", NARANJA),
        ("ROJO  < 40",   "Acción inmediata, reducir carga", ROJO),
    ]
    for i, (lbl, desc, col) in enumerate(legend):
        add_rect(sl, 0.2, 6.5+i*0.22, 2.2, 0.2, col)
        add_text(sl, lbl, 0.22, 6.51+i*0.22, 2.1, 0.2,
                 size=8, bold=True, color=BLANCO)
        add_text(sl, desc, 2.5, 6.51+i*0.22, 6.0, 0.2,
                 size=8, color=NEGRO)

    add_footer(sl)


def slide10_rate_optimo(prs):
    """S10 — Tabla de rates recomendados por régimen."""
    sl = add_slide(prs)
    fill_bg(sl, GRIS_CLARO)
    add_header_bar(sl, "Rate Óptimo por Escenario Operacional",
                   "Tabla prescriptiva validada en backtesting 90 días · SAG1 −1.4% TPH · SAG2 −0.7% TPH")

    img = fig("outputs/figures/06_Rates/optimizacion_rates/07_Matriz_Decision_Operacion.png")
    if img:
        add_img(sl, img, 0.2, 1.25, 7.8, 4.8)

    # Tabla de rates
    add_rect(sl, 8.2, 1.25, 4.9, 0.45, AZUL)
    cols = ["Régimen", "SAG1 Rate", "SAG2 Rate"]
    cwi  = [1.5, 1.65, 1.65]
    x = 8.25
    for c, w in zip(cols, cwi):
        add_text(sl, c, x, 1.3, w, 0.35, size=10, bold=True,
                 color=BLANCO, align=PP_ALIGN.CENTER)
        x += w

    rate_rows = [
        ("AGRESIVO",    "87–105%", "90–105%", VERDE),
        ("NORMAL",      "72–95%",  "82–100%", AMARILLO),
        ("CONSERVADOR", "58–78%",  "76–94%",  NARANJA),
        ("EMERGENCIA",  "50–64%",  "68–82%",  ROJO),
    ]
    for ri, (reg, r1, r2, col) in enumerate(rate_rows):
        bg = GRIS_CLARO if ri % 2 == 0 else BLANCO
        add_rect(sl, 8.2, 1.72+ri*0.52, 4.9, 0.52, bg)
        add_rect(sl, 8.2, 1.72+ri*0.52, 1.5, 0.52, col)
        add_text(sl, reg, 8.2, 1.74+ri*0.52, 1.5, 0.48,
                 size=9.5, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
        add_text(sl, r1, 9.75, 1.74+ri*0.52, 1.65, 0.48,
                 size=10, color=NEGRO, align=PP_ALIGN.CENTER)
        add_text(sl, r2, 11.4, 1.74+ri*0.52, 1.65, 0.48,
                 size=10, color=NEGRO, align=PP_ALIGN.CENTER)

    add_text(sl, "Rate = % del P90 histórico\nSAG1 P90 = 1,454 TPH   ·   SAG2 P90 = 2,516 TPH",
             8.2, 3.85, 4.9, 0.6, size=9, color=AZUL_MED)

    add_text(sl, "Activación del régimen:", 8.2, 4.55, 4.9, 0.35,
             size=10, bold=True, color=AZUL)
    regs_txt = [
        "EMERGENCIA: autonomía < 1.0 h  O  pila < 15%",
        "CONSERVADOR: autonomía < 2.5 h  O  pila < 20%",
        "NORMAL: condición estándar operacional",
        "AGRESIVO: PRE-T8 confirmado, pilas > 40%",
    ]
    for i, t in enumerate(regs_txt):
        add_text(sl, t, 8.2, 4.93+i*0.36, 4.9, 0.34,
                 size=9, color=NEGRO)

    add_footer(sl)


def slide11_kpi_iro(prs):
    """S11 — Índice de Resiliencia Operacional (IRO)."""
    sl = add_slide(prs)
    fill_bg(sl, GRIS_CLARO)
    add_header_bar(sl, "KPI de Resiliencia Operacional — IRO",
                   "Índice compuesto para monitoreo en sala de control")

    img = fig("outputs/figures/11_Modelo_Causal/prescriptivo/P1_IVO_Resiliencia.png")
    if img:
        add_img(sl, img, 0.2, 1.25, 6.5, 5.5)

    # Definición IRO
    add_rect(sl, 6.9, 1.25, 6.2, 5.5, AZUL)
    add_text(sl, "DEFINICIÓN IRO", 7.1, 1.35, 5.9, 0.4,
             size=14, bold=True, color=BLANCO)
    add_text(sl,
             "IRO = f(Inventario, Autonomía,\nRate, Duración T8, Estado Correas)",
             7.1, 1.78, 5.9, 0.7, size=13, color=AMARILLO)

    componentes = [
        ("Inventario (pila)", "Nivel relativo SAG1/SAG2 vs crítico",          "25%"),
        ("Autonomía",         "Horas disponibles antes agotamiento",           "30%"),
        ("Rate operado",      "% vs rate prescrito por régimen",               "20%"),
        ("Duración T8",       "Horas de ventana activa acumuladas",            "15%"),
        ("Estado correas",    "correa_315 disponible = 1, parada = penaliza",  "10%"),
    ]
    add_text(sl, "Componentes y ponderación:", 7.1, 2.55, 5.9, 0.3,
             size=11, bold=True, color=RGBColor(0xB0,0xC4,0xDE))
    for i, (comp, desc, pct) in enumerate(componentes):
        add_rect(sl, 6.9, 2.9+i*0.6, 5.8, 0.55, RGBColor(0x12,0x2C,0x52))
        add_text(sl, comp, 6.95, 2.92+i*0.6, 2.3, 0.48,
                 size=10, bold=True, color=AMARILLO)
        add_text(sl, desc, 9.3, 2.92+i*0.6, 2.9, 0.48,
                 size=8.5, color=BLANCO)
        add_text(sl, pct, 12.25, 2.92+i*0.6, 0.7, 0.48,
                 size=10, bold=True, color=VERDE, align=PP_ALIGN.CENTER)

    add_text(sl, "Escala IRO:  > 80 VERDE · 60-80 AMARILLO · 40-60 NARANJA · < 40 ROJO",
             6.9, 5.95, 6.2, 0.55, size=10, bold=True, color=BLANCO)

    add_footer(sl)


def slide12_simulaciones(prs):
    """S12 — Simulaciones ODE y backtesting."""
    sl = add_slide(prs)
    fill_bg(sl, GRIS_CLARO)
    add_header_bar(sl, "Simulaciones y Backtesting del Sistema RT",
                   "Validación en 90 días históricos: sistema RT vs operador real")

    img1 = fig("outputs/figures/06_Rates/sistema_rt/F03_Backtesting_Pilas.png")
    img2 = fig("outputs/figures/06_Rates/sistema_rt/F05_Metricas_Backtesting.png")
    img3 = fig("outputs/figures/11_Modelo_Causal/decision_operacional/05_Simulacion_T8.png")

    add_img(sl, img1, 0.2, 1.25, 6.3, 3.2)
    add_img(sl, img2, 6.7, 1.25, 6.3, 3.2)
    add_img(sl, img3, 0.2, 4.55, 12.8, 2.25)

    # Resultado numérico
    add_rect(sl, 0.2, 4.5, 12.9, 0.04, AZUL)
    resultados = [
        ("SAG1 TPH", "−1.4%", "✓ dentro ±3%", VERDE),
        ("SAG2 TPH", "−0.7%", "✓ dentro ±3%", VERDE),
        ("SAG1 Agot.", "−1.1%", "≈ límite estructural", AMARILLO),
        ("SAG2 Agot.", "−3.5%", "✓ mejora relevante", VERDE),
    ]
    for i, (lbl, val, nota, col) in enumerate(resultados):
        x = 0.2 + i * 3.22
        add_rect(sl, x, 6.85, 3.1, 0.45, col)
        add_text(sl, f"{lbl}: {val}  ({nota})", x+0.05, 6.87, 3.0, 0.42,
                 size=9, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)

    add_footer(sl)


def slide13_recomendaciones(prs):
    """S13 — Recomendaciones operacionales concretas."""
    sl = add_slide(prs)
    fill_bg(sl, GRIS_CLARO)
    add_header_bar(sl, "Recomendaciones Operacionales",
                   "Acciones concretas clasificadas por horizonte y responsable")

    add_rect(sl, 0.2, 1.2, 4.15, 5.8, RGBColor(0xD4,0xEB,0xDA))   # verde claro
    add_rect(sl, 4.6, 1.2, 4.0, 5.8, RGBColor(0xFE,0xF3,0xCD))    # amarillo claro
    add_rect(sl, 8.8, 1.2, 4.3, 5.8, RGBColor(0xFC,0xE4,0xD6))    # naranja claro

    add_rect(sl, 0.2, 1.2, 4.15, 0.45, VERDE)
    add_rect(sl, 4.6, 1.2, 4.0, 0.45, AMARILLO)
    add_rect(sl, 8.8, 1.2, 4.3, 0.45, NARANJA)

    add_text(sl, "INMEDIATO (0–30 días)", 0.25, 1.22, 4.0, 0.4,
             size=11, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    add_text(sl, "CORTO PLAZO (1–3 meses)", 4.65, 1.22, 3.85, 0.4,
             size=11, bold=True, color=NEGRO, align=PP_ALIGN.CENTER)
    add_text(sl, "MEDIANO PLAZO (3–6 meses)", 8.85, 1.22, 4.1, 0.4,
             size=11, bold=True, color=NEGRO, align=PP_ALIGN.CENTER)

    inm = [
        "Ajustar alarmas pila\nSAG1/SAG2 a 18%\n(reemplaza 70%)",
        "Implementar tabla\nrates prescritos\nen sala de control",
        "Activar protocolo\nPRE-T8: subir rates\n4 h antes confirmado",
        "Registrar hora inicio\nT8 en sistema SCADA\n(trazabilidad)",
    ]
    cp = [
        "Integrar semáforo IRO\nen Power BI\n(sala de control)",
        "Capacitar operadores\nen nuevos umbrales\ny protocolos",
        "Validar parámetros\npila con Ingeniería\nde Proceso",
        "Pipeline mensual\nPAM para reentrenar\nmodelos",
    ]
    mp = [
        "Plan de mejora\ncorrea 315\n(restricción crítica)",
        "Dashboard autonomía\ntime-series para\nJefe Turno",
        "Extender modelo\na PMC y UNITARIO\ncon misma lógica",
        "Integración SCADA\nreal-time para\nalertas automáticas",
    ]
    for col_items, x in [(inm, 0.25), (cp, 4.65), (mp, 8.85)]:
        for i, t in enumerate(col_items):
            add_rect(sl, x, 1.72+i*1.15, 4.0 if x<8 else 4.2, 1.08,
                     BLANCO)
            add_text(sl, t, x+0.1, 1.75+i*1.15, 3.8 if x<8 else 4.0, 1.03,
                     size=9.5, color=NEGRO)

    add_footer(sl)


def slide14_roadmap(prs):
    """S14 — Roadmap de implementación."""
    sl = add_slide(prs)
    fill_bg(sl, GRIS_CLARO)
    add_header_bar(sl, "Roadmap de Implementación",
                   "Hoja de ruta en 3 fases hacia sistema prescriptivo en producción")

    phases = [
        ("FASE 1\nJulio 2026",
         ["Ajustar alarmas pila (18%)",
          "Distribuir tabla rates prescritos",
          "Capacitación operadores (4 h)",
          "Protocolo PRE-T8 formal"],
         VERDE, 0.2),
        ("FASE 2\nAgo–Sep 2026",
         ["Dashboard IRO Power BI",
          "Pipeline PAM mensual",
          "Validación parámetros con IP",
          "Extensión PMC/UNITARIO"],
         AMARILLO, 4.55),
        ("FASE 3\nOct–Dic 2026",
         ["Plan mejora correa 315",
          "Integración SCADA RT",
          "Alertas automáticas turno",
          "Revisión anual del modelo"],
         NARANJA, 8.9),
    ]
    for label, items, col, x in phases:
        add_rect(sl, x, 1.25, 4.1, 0.65, col)
        add_text(sl, label, x+0.1, 1.27, 3.9, 0.62,
                 size=13, bold=True, color=BLANCO if col != AMARILLO else NEGRO,
                 align=PP_ALIGN.CENTER)
        add_rect(sl, x, 1.92, 4.1, 4.7, BLANCO)
        for i, it in enumerate(items):
            add_rect(sl, x+0.1, 2.0+i*1.1, 0.08, 0.45, col)
            add_text(sl, it, x+0.3, 2.0+i*1.1, 3.7, 0.45,
                     size=11, color=NEGRO)

    # Hitos transversales
    add_rect(sl, 0.2, 6.7, 12.9, 0.6, AZUL)
    add_text(sl,
             "Hitos transversales: KPI IRO en dashboard desde Fase 1  ·  "
             "Revisión trimestral Jefe Planta  ·  "
             "Entregable PAM cada ciclo planificación",
             0.35, 6.75, 12.5, 0.5, size=10.5, color=BLANCO)

    add_footer(sl)


def slide15_conclusiones(prs):
    """S15 — Conclusiones y próximos pasos."""
    sl = add_slide(prs)
    fill_bg(sl, AZUL)

    add_rect(sl, 0, 0, 13.33, 0.12, AMARILLO)
    add_text(sl, "Conclusiones y Próximos Pasos", 0.5, 0.2, 12, 0.75,
             size=28, bold=True, color=BLANCO, align=PP_ALIGN.LEFT)
    add_text(sl, "División El Teniente · Analítica Prescriptiva Molienda SAG · Junio 2026",
             0.5, 0.88, 12, 0.4, size=13,
             color=RGBColor(0xB0,0xC4,0xDE), align=PP_ALIGN.LEFT)

    add_rect(sl, 0.2, 1.38, 12.9, 0.04, AMARILLO)

    conclusiones = [
        ("¿Cuál es el principal cuello de botella?",
         "Correa 315 parada 49% del tiempo → SAG1 en déficit crónico. "
         "Ninguna optimización de rates puede compensar esta restricción estructural."),
        ("¿Cuándo hay riesgo real?",
         "Pila SAG1 o SAG2 < 18%, o autonomía < 2.5 h. "
         "El umbral operacional histórico de 70% es excesivamente conservador."),
        ("¿Qué hacer antes de una T8?",
         "Activar régimen AGRESIVO 4 h antes si pilas > 40%. "
         "SAG2 puede pre-cargar independientemente de SAG1."),
        ("¿Cuánto mejora el sistema RT?",
         "SAG1 −1.4% TPH, SAG2 −0.7% TPH (dentro del ±3% objetivo). "
         "Reducción agotamiento SAG2 −3.5 pp. SAG1 limitado por correa 315."),
        ("¿Cuál es el primer paso concreto?",
         "Esta semana: ajustar alarma pila de 70% a 18% y distribuir tabla de rates prescritos."),
    ]

    for i, (preg, resp) in enumerate(conclusiones):
        y = 1.5 + i * 1.12
        add_rect(sl, 0.2, y, 12.9, 1.05, RGBColor(0x12,0x28,0x4F))
        add_rect(sl, 0.2, y, 0.1, 1.05, AMARILLO)
        add_text(sl, preg, 0.4, y+0.04, 12.4, 0.37,
                 size=11, bold=True, color=AMARILLO)
        add_text(sl, resp, 0.4, y+0.42, 12.4, 0.57,
                 size=10.5, color=BLANCO)

    add_rect(sl, 0, 7.12, 13.33, 0.38, RGBColor(0x10, 0x1E, 0x3A))
    add_text(sl, "División El Teniente · Analítica & Datos · CONFIDENCIAL",
             0.3, 7.15, 10, 0.3, size=9,
             color=RGBColor(0x70,0x80,0x90))


# ─────────────────────────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    print("Generando presentación — 15 slides...")

    slide01_portada(prs)                 ; print("  S01 Portada")
    slide02_resumen_ejecutivo(prs)       ; print("  S02 Resumen Ejecutivo")
    slide03_problema_operacional(prs)    ; print("  S03 Problema Operacional")
    slide04_cadena_causal(prs)           ; print("  S04 Cadena Causal")
    slide05_hallazgos_clave(prs)         ; print("  S05 Hallazgos Clave")
    slide06_impacto_t8(prs)              ; print("  S06 Impacto T8")
    slide07_autonomia_pilas(prs)         ; print("  S07 Autonomía Pilas")
    slide08_validacion_reglas(prs)       ; print("  S08 Validación Reglas")
    slide09_matriz_riesgo(prs)           ; print("  S09 Matriz Riesgo")
    slide10_rate_optimo(prs)             ; print("  S10 Rate Óptimo")
    slide11_kpi_iro(prs)                 ; print("  S11 KPI IRO")
    slide12_simulaciones(prs)            ; print("  S12 Simulaciones")
    slide13_recomendaciones(prs)         ; print("  S13 Recomendaciones")
    slide14_roadmap(prs)                 ; print("  S14 Roadmap")
    slide15_conclusiones(prs)            ; print("  S15 Conclusiones")

    out_dir = ROOT / "outputs" / "presentations"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "20260625_Analitica_Prescriptiva_Molienda_T8.pptx"
    prs.save(str(out_path))
    print(f"\nOK Guardado: {out_path}")
    print(f"  Tamanio: {out_path.stat().st_size / 1024:.0f} KB")

if __name__ == "__main__":
    main()
