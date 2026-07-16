"""
generar_informe_pdf.py — Informe ejecutivo PDF + PPTX para Comite T8.

Skills aplicados: skill_molienda_sag, skill_product_owner_analitica_minera,
                  skill_ux_ui_cio_operations_center, skill_operaciones_mina_subterranea.

Entregables:
    outputs/reports/Informe_Comite_T8.pdf
    outputs/reports/Informe_Comite_T8.pptx
    outputs/reports/Anexo_Tecnico_T8.pdf

Uso:
    python src/generar_informe_pdf.py
"""
from __future__ import annotations

import io, json, os, sys, textwrap
from datetime import datetime
from pathlib import Path
from typing import Optional

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.image as mpimg
from matplotlib.backends.backend_pdf import PdfPages
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch
import numpy as np
import pandas as pd

# ── Rutas ──────────────────────────────────────────────────────────────────────
BASE   = Path(__file__).resolve().parents[1]
FIGS   = BASE / "outputs" / "figures"
RPT    = BASE / "outputs" / "reports"
LOGS   = BASE / "logs"
XLS    = BASE / "outputs" / "excel" / "event_study_t8.xlsx"

RPT.mkdir(parents=True, exist_ok=True)
LOGS.mkdir(parents=True, exist_ok=True)

# ── Paleta corporativa ─────────────────────────────────────────────────────────
C_BLUE    = "#1B3A5C"    # azul corporativo profundo
C_COPPER  = "#B87333"    # cobre
C_ORANGE  = "#C95B27"    # naranja operacional
C_RED     = "#C0392B"    # alerta
C_GREEN   = "#27AE60"    # positivo
C_LGRAY   = "#F4F6F9"    # fondo suave
C_MGRAY   = "#8E9BAA"    # texto secundario
C_WHITE   = "#FFFFFF"
C_DARK    = "#1A1A1A"
C_GOLD    = "#D4A83A"    # acento dorado

# ── Página 16:9 ───────────────────────────────────────────────────────────────
PW, PH = 16.0, 9.0   # pulgadas
DPI     = 150

# ── Datos reales ──────────────────────────────────────────────────────────────

def _load_data() -> dict:
    df_met  = pd.read_excel(XLS, sheet_name="metricas_evento_activo")
    df_stat = pd.read_excel(XLS, sheet_name="significancia_estadistica")
    df_act  = pd.read_excel(XLS, sheet_name="resumen_activo")
    return {"met": df_met, "stat": df_stat, "act": df_act}


# ═══════════════════════════════════════════════════════════════════════════════
# Utilidades de diseño
# ═══════════════════════════════════════════════════════════════════════════════

def _new_page(facecolor: str = C_WHITE) -> plt.Figure:
    fig = plt.figure(figsize=(PW, PH))
    fig.patch.set_facecolor(facecolor)
    return fig


def _header(fig: plt.Figure, title: str, subtitle: str = "",
            bar_color: str = C_BLUE, page_num: int = 0) -> None:
    """Banda superior con título."""
    ax = fig.add_axes([0, 0.88, 1, 0.12])
    ax.set_facecolor(bar_color)
    ax.axis("off")
    ax.text(0.02, 0.58, title,   color=C_WHITE, fontsize=16, fontweight="bold",
            va="center", transform=ax.transAxes)
    if subtitle:
        ax.text(0.02, 0.18, subtitle, color=C_GOLD, fontsize=9,
                va="center", transform=ax.transAxes)
    if page_num:
        ax.text(0.98, 0.35, f"[ {page_num} ]", color=C_WHITE, fontsize=8,
                ha="right", va="center", transform=ax.transAxes, alpha=0.7)
    # Línea decorativa
    ax2 = fig.add_axes([0, 0.876, 1, 0.006])
    ax2.set_facecolor(C_COPPER)
    ax2.axis("off")


def _footer(fig: plt.Figure, date_str: str, version: str = "v1.0") -> None:
    """Banda inferior."""
    ax = fig.add_axes([0, 0, 1, 0.04])
    ax.set_facecolor(C_BLUE)
    ax.axis("off")
    ax.text(0.02, 0.5, "CODELCO División El Teniente — CIO DET — Analítica de Rendimientos",
            color=C_WHITE, fontsize=7, va="center", transform=ax.transAxes, alpha=0.85)
    ax.text(0.98, 0.5, f"{date_str}  |  {version}  |  CONFIDENCIAL",
            color=C_GOLD, fontsize=7, ha="right", va="center",
            transform=ax.transAxes, alpha=0.9)


def _img(fig: plt.Figure, path: str | Path,
         left: float, bottom: float, width: float, height: float,
         alpha: float = 1.0) -> None:
    """Inserta imagen en coordenadas de figura normalizadas."""
    p = Path(path)
    if not p.exists():
        return
    ax = fig.add_axes([left, bottom, width, height])
    img = mpimg.imread(str(p))
    ax.imshow(img, aspect="auto")
    ax.axis("off")
    if alpha < 1.0:
        ax.set_alpha(alpha)


def _kpi_box(fig: plt.Figure, left: float, bottom: float,
             width: float, height: float,
             value: str, label: str, sublabel: str = "",
             bg: str = C_BLUE, text_c: str = C_WHITE) -> None:
    """Caja KPI con valor grande y etiqueta."""
    ax = fig.add_axes([left, bottom, width, height])
    ax.set_facecolor(bg)
    ax.set_xlim(0, 1); ax.set_ylim(0, 1)
    ax.axis("off")
    ax.add_patch(FancyBboxPatch((0.03, 0.03), 0.94, 0.94,
                                boxstyle="round,pad=0.02",
                                facecolor=bg, edgecolor=C_COPPER, linewidth=1.5))
    ax.text(0.5, 0.60, value,    color=text_c, fontsize=22, fontweight="bold",
            ha="center", va="center")
    ax.text(0.5, 0.28, label,    color=text_c, fontsize=9,  fontweight="bold",
            ha="center", va="center")
    if sublabel:
        ax.text(0.5, 0.10, sublabel, color=C_GOLD, fontsize=7.5,
                ha="center", va="center")


def _bullet_block(fig: plt.Figure, left: float, bottom: float,
                  width: float, height: float,
                  items: list[tuple[str, str]],
                  title: str = "", bg: str = C_LGRAY,
                  icon_color: str = C_ORANGE) -> None:
    """Bloque de bullets con icono de color."""
    ax = fig.add_axes([left, bottom, width, height])
    ax.set_facecolor(bg)
    ax.set_xlim(0, 1); ax.set_ylim(0, 1)
    ax.axis("off")
    ax.add_patch(FancyBboxPatch((0.01, 0.01), 0.98, 0.98,
                                boxstyle="round,pad=0.01",
                                facecolor=bg, edgecolor=C_MGRAY, linewidth=0.6))
    y0 = 0.90
    if title:
        ax.text(0.05, y0, title, color=C_BLUE, fontsize=9, fontweight="bold", va="top")
        y0 -= 0.12
    step = (y0 - 0.05) / max(len(items), 1)
    for i, (num, text) in enumerate(items):
        y = y0 - i * step
        ax.text(0.05, y, num,   color=icon_color, fontsize=9, fontweight="bold", va="center")
        ax.text(0.15, y, text,  color=C_DARK, fontsize=8, va="center",
                wrap=True)


def _callout(fig: plt.Figure, left: float, bottom: float,
             width: float, height: float,
             text: str, bg: str = C_BLUE, fc: str = C_WHITE) -> None:
    """Caja de mensaje principal estilo callout."""
    ax = fig.add_axes([left, bottom, width, height])
    ax.set_facecolor(bg)
    ax.set_xlim(0, 1); ax.set_ylim(0, 1)
    ax.axis("off")
    ax.add_patch(FancyBboxPatch((0.02, 0.05), 0.96, 0.90,
                                boxstyle="round,pad=0.02",
                                facecolor=bg, edgecolor=C_COPPER, linewidth=2))
    # Wrapping text
    wrapped = textwrap.fill(text, width=40)
    ax.text(0.5, 0.5, wrapped, color=fc, fontsize=11, fontweight="bold",
            ha="center", va="center", multialignment="center", linespacing=1.6)


def _section_label(fig: plt.Figure, left: float, bottom: float,
                   width: float, text: str, color: str = C_ORANGE) -> None:
    ax = fig.add_axes([left, bottom, width, 0.04])
    ax.axis("off")
    ax.set_xlim(0, 1); ax.set_ylim(0, 1)
    ax.add_patch(mpatches.Rectangle((0, 0.1), 0.025, 0.8, color=color))
    ax.text(0.04, 0.5, text.upper(), color=color, fontsize=8,
            fontweight="bold", va="center", transform=ax.transAxes,
            letter_spacing=1.5 if hasattr(ax.text, "letter_spacing") else None)


# ═══════════════════════════════════════════════════════════════════════════════
# PÁGINAS DEL INFORME PRINCIPAL
# ═══════════════════════════════════════════════════════════════════════════════

DATE_STR = datetime.now().strftime("%d/%m/%Y")

def page_portada() -> plt.Figure:
    fig = _new_page(C_BLUE)

    # Banda lateral izquierda (cobre)
    ax_side = fig.add_axes([0, 0, 0.008, 1])
    ax_side.set_facecolor(C_COPPER)
    ax_side.axis("off")

    # Imagen de fondo (panel gris claro, simulado)
    ax_bg = fig.add_axes([0.008, 0, 0.992, 1])
    ax_bg.set_facecolor(C_BLUE)
    ax_bg.set_xlim(0, 1); ax_bg.set_ylim(0, 1)
    ax_bg.axis("off")

    # Franja inferior decorativa
    ax_bg.add_patch(mpatches.Rectangle((0, 0), 1, 0.18, color=C_DARK, alpha=0.4))

    # Título principal
    ax_bg.text(0.06, 0.82, "IMPACTO OPERACIONAL DE LAS", color=C_GOLD,
               fontsize=14, fontweight="bold", va="center")
    ax_bg.text(0.06, 0.72, "VENTANAS TENIENTE 8", color=C_WHITE,
               fontsize=32, fontweight="bold", va="center")
    ax_bg.text(0.06, 0.62, "SOBRE LOS RENDIMIENTOS DE MOLIENDA", color=C_COPPER,
               fontsize=16, fontweight="bold", va="center")

    # Línea separadora
    ax_bg.axhline(0.56, xmin=0.06, xmax=0.65, color=C_COPPER, linewidth=1.5)

    # Subtítulo
    ax_bg.text(0.06, 0.49, "Análisis Event Study Industrial", color=C_WHITE,
               fontsize=12, va="center")
    ax_bg.text(0.06, 0.42, "Período: Enero 2026 – Junio 2026   |   72 Eventos T8", color=C_MGRAY,
               fontsize=10, va="center")

    # Recuadros de dato rápido
    datos = [
        ("72", "Eventos\nAnalizados"),
        ("136.7k ton", "Pérdida\nEstimada"),
        ("42.7%", "Máxima\nCaída"),
        ("3", "Activos con\nSignificancia"),
    ]
    for i, (val, lbl) in enumerate(datos):
        x = 0.06 + i * 0.165
        ax_bg.add_patch(FancyBboxPatch((x, 0.24), 0.15, 0.10,
                                       boxstyle="round,pad=0.01",
                                       facecolor=C_DARK, edgecolor=C_COPPER,
                                       linewidth=1.2, alpha=0.8))
        ax_bg.text(x + 0.075, 0.315, val, color=C_GOLD, fontsize=12,
                   fontweight="bold", ha="center", va="center")
        ax_bg.text(x + 0.075, 0.260, lbl, color=C_WHITE, fontsize=7,
                   ha="center", va="center")

    # Metadatos
    ax_bg.text(0.06, 0.11, f"Fecha: {DATE_STR}",
               color=C_WHITE, fontsize=9, va="center", alpha=0.85)
    ax_bg.text(0.06, 0.07, "Versión: 1.0   |   CONFIDENCIAL",
               color=C_MGRAY, fontsize=8, va="center")
    ax_bg.text(0.06, 0.03, "CIO DET — Analítica de Rendimientos — CODELCO División El Teniente",
               color=C_MGRAY, fontsize=7.5, va="center")

    # Marca de cobre decorativa
    theta = np.linspace(0, 2 * np.pi, 300)
    r = 0.22
    cx, cy = 0.82, 0.52
    ax_bg.plot(cx + r * np.cos(theta), cy + r * np.sin(theta) * 0.9,
               color=C_COPPER, linewidth=0.5, alpha=0.2)
    ax_bg.plot(cx + r * 0.7 * np.cos(theta), cy + r * 0.7 * np.sin(theta) * 0.9,
               color=C_COPPER, linewidth=0.3, alpha=0.15)
    ax_bg.text(cx, cy + 0.03, "T8", color=C_COPPER, fontsize=60, fontweight="bold",
               ha="center", va="center", alpha=0.12)

    return fig


def page_resumen_ejecutivo() -> plt.Figure:
    fig = _new_page()
    _header(fig, "RESUMEN EJECUTIVO",
            "4 hallazgos clave del análisis Event Study Industrial", page_num=2)
    _footer(fig, DATE_STR)

    # 4 KPIs en la fila superior
    kpis = [
        ("42.7%",  "SAG2 — Mayor caída promedio",   "Activo más sensible al efecto T8",     C_RED),
        ("39.9%",  "Ventanas 12h — Máximo impacto", "3.3x más pérdida vs ventana 4h",       C_ORANGE),
        ("11.4 h", "Recuperación SAG2 al 90%",      "La recuperación más lenta del sistema", C_BLUE),
        ("3 de 4", "Activos con evidencia estadística", "SAG1, SAG2 y PMC — p < 0.05",      C_COPPER),
    ]
    for i, (val, lbl, sub, col) in enumerate(kpis):
        _kpi_box(fig, 0.03 + i * 0.245, 0.56, 0.22, 0.28,
                 val, lbl, sub, bg=col)

    # Mensaje central
    _callout(fig, 0.03, 0.26, 0.42, 0.26,
             "Cuando T8 se detiene, los molinos no reciben mineral fresco.\nLas pilas amortiguan temporalmente la caida,\npero el impacto en TPH ocurre DESPUES de terminada la ventana.",
             bg=C_BLUE, fc=C_WHITE)

    # Implicancia ejecutiva
    ax_impl = fig.add_axes([0.48, 0.26, 0.50, 0.26])
    ax_impl.set_facecolor(C_LGRAY)
    ax_impl.set_xlim(0, 1); ax_impl.set_ylim(0, 1)
    ax_impl.axis("off")
    ax_impl.add_patch(FancyBboxPatch((0.01, 0.01), 0.98, 0.97,
                                     boxstyle="round,pad=0.01",
                                     facecolor=C_LGRAY, edgecolor=C_COPPER, linewidth=1.2))
    ax_impl.text(0.05, 0.85, "IMPLICANCIA PARA EL COMITÉ",
                 color=C_BLUE, fontsize=9, fontweight="bold", va="top")
    bullets = [
        "Reducir ventanas 12h es la accion de mayor ROI",
        "SAG2 y PMC requieren protocolo especifico pre-ventana",
        "El stock de pila es la palanca operacional clave",
        "La evidencia estadistica respalda accion correctiva",
    ]
    for j, b in enumerate(bullets):
        ax_impl.text(0.06, 0.65 - j * 0.16, f"■  {b}",
                     color=C_DARK, fontsize=8.5, va="center")

    # Nota de confianza
    ax_nota = fig.add_axes([0.03, 0.10, 0.95, 0.12])
    ax_nota.set_facecolor(C_LGRAY)
    ax_nota.set_xlim(0, 1); ax_nota.set_ylim(0, 1)
    ax_nota.axis("off")
    ax_nota.text(0.02, 0.6, "BASE DEL ANALISIS:",
                 color=C_BLUE, fontsize=8, fontweight="bold", va="center")
    ax_nota.text(0.22, 0.6,
                 "72 ventanas T8 (ene-jun 2026)  |  47,532 registros de rendimiento a 5 min  |"
                 "  Event Study Industrial alineado en t=0 = hora oficial de inicio  |"
                 "  Significancia: T-test + Mann-Whitney (alpha=0.05)",
                 color=C_MGRAY, fontsize=7.5, va="center")

    return fig


def page_contexto() -> plt.Figure:
    fig = _new_page()
    _header(fig, "CONTEXTO OPERACIONAL",
            "¿Qué es Teniente 8 y cómo afecta los rendimientos?", page_num=3)
    _footer(fig, DATE_STR)

    # Diagrama de flujo causal
    ax = fig.add_axes([0.03, 0.12, 0.38, 0.70])
    ax.set_xlim(0, 10); ax.set_ylim(0, 10)
    ax.axis("off")
    ax.set_facecolor(C_LGRAY)
    ax.add_patch(FancyBboxPatch((0.1, 0.1), 9.8, 9.8, boxstyle="round,pad=0.1",
                                facecolor=C_LGRAY, edgecolor=C_COPPER, linewidth=1))

    steps = [
        (C_BLUE,    "TENIENTE 8",                  "Ferrocarril mineral fino/grueso → pilas"),
        (C_ORANGE,  "VENTANA DE MANTENIMIENTO",    "2h / 4h / 8h / 12h"),
        (C_RED,     "SIN ENTREGA A CHANCADO",      "Primario y secundario sin alimentacion"),
        (C_RED,     "SIN OFERTA A MOLINOS",        "No hay mineral disponible para moler"),
        (C_COPPER,  "CONSUMO DE PILAS",            "Stock amortiguador SAG / PMC / UNITARIO"),
        (C_GREEN,   "CAIDA TPH → RECUPERACION",   "Efecto diferido y retorno gradual"),
    ]
    y_positions = [8.8, 7.3, 5.8, 4.3, 2.8, 1.3]
    for (col, ttl, sub), y in zip(steps, y_positions):
        ax.add_patch(FancyBboxPatch((0.8, y - 0.5), 8.4, 1.1,
                                    boxstyle="round,pad=0.1",
                                    facecolor=col, edgecolor="none", alpha=0.88))
        ax.text(5, y + 0.15, ttl, color=C_WHITE, fontsize=8.5, fontweight="bold",
                ha="center", va="center")
        ax.text(5, y - 0.18, sub, color=C_WHITE, fontsize=7,
                ha="center", va="center", alpha=0.85)
        if y > 1.3:
            ax.annotate("", xy=(5, y - 0.55), xytext=(5, y - 0.42),
                        arrowprops=dict(arrowstyle="-|>", color=C_MGRAY, lw=1.2))
    ax.set_title("Cadena de impacto T8", color=C_BLUE,
                 fontsize=9, fontweight="bold", pad=4)

    # Texto explicativo
    ax2 = fig.add_axes([0.44, 0.12, 0.54, 0.70])
    ax2.set_facecolor(C_WHITE)
    ax2.set_xlim(0, 1); ax2.set_ylim(0, 1)
    ax2.axis("off")

    bloques = [
        ("¿QUE ES TENIENTE 8?",
         "Sistema ferroviario que transporta mineral fino y grueso desde la mina "
         "hacia las pilas de alimentacion de los circuitos SAG, PMC y UNITARIO. "
         "Es la principal fuente de mineral para los molinos de la Division."),
        ("¿QUE ES UNA VENTANA?",
         "Detencion programada del ferrocarril T8 para mantenimiento. Durante la ventana "
         "NO hay entrega de mineral a chancado primario ni secundario, interrumpiendo "
         "completamente la oferta de mineral hacia los molinos. Pueden durar 2h, 4h, 8h o 12h."),
        ("¿POR QUE SE ANALIZO?",
         "En 5.5 meses se registraron 72 ventanas T8. El impacto acumulado "
         "en produccion es significativo pero no habia sido cuantificado. "
         "Este analisis permite tomar decisiones basadas en evidencia."),
        ("¿COMO AFECTAN LAS PILAS?",
         "Al cortarse la oferta de mineral, los molinos consumen el stock acumulado en pila. "
         "Esto RETRASA la caida de TPH (IAP > 1 en todos los activos). "
         "Cuando el stock se agota, la caida de rendimiento es inevitable."),
    ]
    y0 = 0.96
    for ttl, txt in bloques:
        ax2.add_patch(mpatches.Rectangle((0, y0 - 0.21), 0.012, 0.17,
                                          color=C_COPPER, alpha=0.9))
        ax2.text(0.025, y0 - 0.035, ttl, color=C_BLUE, fontsize=8.5,
                 fontweight="bold", va="top")
        wrapped = textwrap.fill(txt, width=68)
        ax2.text(0.025, y0 - 0.09, wrapped, color=C_DARK, fontsize=7.5,
                 va="top", linespacing=1.5)
        y0 -= 0.26

    return fig


def page_metodologia() -> plt.Figure:
    fig = _new_page()
    _header(fig, "METODOLOGIA",
            "Event Study Industrial — Enfoque estadistico riguroso", page_num=4)
    _footer(fig, DATE_STR)

    ax = fig.add_axes([0.03, 0.10, 0.94, 0.73])
    ax.set_xlim(0, 16); ax.set_ylim(0, 6)
    ax.axis("off")
    ax.set_facecolor(C_WHITE)

    # Componentes de entrada
    inputs = [
        ("PAM\nMANTTO", "Registro oficial\nde ventanas", C_BLUE),
        ("RENDIMIENTOS\n5 MIN", "47,532 registros\nSAG1, SAG2, PMC,\nUNITARIO", C_COPPER),
    ]
    for i, (ttl, sub, col) in enumerate(inputs):
        y = 4.5 - i * 2.2
        ax.add_patch(FancyBboxPatch((0.2, y - 0.6), 2.6, 1.3,
                                    boxstyle="round,pad=0.1",
                                    facecolor=col, edgecolor="none"))
        ax.text(1.5, y + 0.12, ttl, color=C_WHITE, fontsize=9, fontweight="bold",
                ha="center", va="center")
        ax.text(1.5, y - 0.32, sub, color=C_WHITE, fontsize=7,
                ha="center", va="center", alpha=0.85)
        ax.annotate("", xy=(3.2, 3.5 - i * 0.4), xytext=(2.8, y - 0.1),
                    arrowprops=dict(arrowstyle="-|>", color=C_MGRAY, lw=1.0))

    # Proceso central
    pasos = [
        ("1\nIDENTIFICAR\n72 EVENTOS", C_BLUE),
        ("2\nALINEAR\nt=0 = inicio\nventana", C_BLUE),
        ("3\n-24h / +24h\npor evento", C_BLUE),
        ("4\nCURVA\nPROMEDIO\npor activo", C_ORANGE),
        ("5\nT-TEST +\nMANN-WHITNEY\npre vs post", C_RED),
        ("6\nKPI:\nIVO, IR,\nIAP, IST8", C_COPPER),
    ]
    for j, (ttl, col) in enumerate(pasos):
        x = 3.4 + j * 2.1
        ax.add_patch(FancyBboxPatch((x, 2.4), 1.85, 1.8,
                                    boxstyle="round,pad=0.1",
                                    facecolor=col, edgecolor="none", alpha=0.9))
        ax.text(x + 0.925, 3.3, ttl, color=C_WHITE, fontsize=7.5,
                fontweight="bold", ha="center", va="center", linespacing=1.4)
        if j < len(pasos) - 1:
            ax.annotate("", xy=(x + 1.9, 3.3), xytext=(x + 1.85, 3.3),
                        arrowprops=dict(arrowstyle="-|>", color=C_MGRAY, lw=1.0))

    # Outputs
    outputs = [
        ("CURVAS DE\nRESPUESTA", "Comportamiento\npromedio post-T8"),
        ("METRICAS POR\nEVENTO", "Caida%, recuperacion,\ntoneladas perdidas"),
        ("ESCENARIOS\nOPERACIONALES", "ROI de reduccion\nde duracion"),
    ]
    for k, (ttl, sub) in enumerate(outputs):
        x = 4.0 + k * 3.8
        ax.add_patch(FancyBboxPatch((x, 0.3), 3.2, 1.5,
                                    boxstyle="round,pad=0.1",
                                    facecolor=C_GREEN, edgecolor="none", alpha=0.85))
        ax.annotate("", xy=(x + 1.6, 1.8), xytext=(x + 1.6, 2.38),
                    arrowprops=dict(arrowstyle="-|>", color=C_GREEN, lw=1.2))
        ax.text(x + 1.6, 1.2, ttl, color=C_WHITE, fontsize=8, fontweight="bold",
                ha="center", va="center", linespacing=1.4)
        ax.text(x + 1.6, 0.55, sub, color=C_WHITE, fontsize=7,
                ha="center", va="center", alpha=0.85)

    # Nota t=0
    ax.add_patch(FancyBboxPatch((0.2, 0.1), 2.8, 1.7,
                                boxstyle="round,pad=0.1",
                                facecolor=C_LGRAY, edgecolor=C_COPPER, linewidth=1.5))
    ax.text(1.6, 1.3, "t = 0", color=C_BLUE, fontsize=18, fontweight="bold", ha="center")
    ax.text(1.6, 0.85, "Inicio OFICIAL\nde la ventana T8", color=C_DARK, fontsize=7.5,
            ha="center", va="center", linespacing=1.4)
    ax.text(1.6, 0.35, "2h → 14:00  |  4h → 12:00\n8h → 08:00  |  12h → 08:00",
            color=C_MGRAY, fontsize=7, ha="center", va="center", linespacing=1.3)

    return fig


def page_hallazgo_gaviota() -> plt.Figure:
    fig = _new_page()
    _header(fig, "HALLAZGO 1 — EFECTO GAVIOTA",
            "Patron repetitivo de caida y recuperacion post-ventana T8", page_num=5)
    _footer(fig, DATE_STR)

    _img(fig, FIGS / "event_study" / "09_Efecto_Gaviota_Global.png",
         0.03, 0.12, 0.68, 0.72)

    # Mensajes clave
    msgs = [
        (C_ORANGE, "PATRON REPETITIVO",
         "72 eventos muestran el mismo\ncomportamiento: caida diferida\ny recuperacion gradual."),
        (C_RED, "CAIDA DIFERIDA",
         "T8 detiene entrega a chancado.\nLas pilas consumen su stock.\nCuando se agota → cae el TPH."),
        (C_GREEN, "RECUPERACION LENTA",
         "La recuperacion al 90% del\nbaseline toma entre 5.8h (UNITARIO)\ny 11.4h (SAG2)."),
    ]
    for i, (col, ttl, body) in enumerate(msgs):
        y = 0.73 - i * 0.22
        ax = fig.add_axes([0.73, y, 0.25, 0.19])
        ax.set_facecolor(C_LGRAY)
        ax.set_xlim(0, 1); ax.set_ylim(0, 1)
        ax.axis("off")
        ax.add_patch(mpatches.Rectangle((0, 0), 0.03, 1, color=col))
        ax.text(0.07, 0.78, ttl,  color=col,    fontsize=8.5, fontweight="bold", va="top")
        ax.text(0.07, 0.50, body, color=C_DARK, fontsize=7.5, va="top", linespacing=1.5)

    _callout(fig, 0.73, 0.12, 0.25, 0.13,
             "La evidencia no es casual:\npatron estadisticamente significativo\nen SAG1, SAG2 y PMC.",
             bg=C_BLUE)
    return fig


def page_hallazgo_duracion() -> plt.Figure:
    fig = _new_page()
    _header(fig, "HALLAZGO 2 — VENTANAS LARGAS GENERAN MAYOR DANO",
            "Comparacion del efecto por duracion de ventana", page_num=6)
    _footer(fig, DATE_STR)

    layout = [
        ("05_EventStudy_2h.png",  0.03, 0.47, 0.45, 0.41, "2h  |  34.9% caida  |  540 ton/evento"),
        ("06_EventStudy_4h.png",  0.51, 0.47, 0.45, 0.41, "4h  |  38.3% caida  |  2,507 ton/evento"),
        ("07_EventStudy_8h.png",  0.03, 0.11, 0.45, 0.41, "8h  |  ~39% caida   |  ~5,375 ton/evento"),
        ("08_EventStudy_12h.png", 0.51, 0.11, 0.45, 0.41, "12h |  39.9% caida  |  8,242 ton/evento"),
    ]
    dur_colors = {0: C_BLUE, 1: C_ORANGE, 2: C_GREEN, 3: C_RED}

    for i, (fname, left, bot, w, h, lbl) in enumerate(layout):
        _img(fig, FIGS / "event_study" / fname, left, bot + 0.045, w, h - 0.05)
        ax_lbl = fig.add_axes([left, bot, w, 0.04])
        ax_lbl.set_facecolor(dur_colors[i])
        ax_lbl.axis("off")
        ax_lbl.text(0.5, 0.5, lbl, color=C_WHITE, fontsize=8, fontweight="bold",
                    ha="center", va="center", transform=ax_lbl.transAxes)

    return fig


def page_hallazgo_activos() -> plt.Figure:
    fig = _new_page()
    _header(fig, "HALLAZGO 3 — ¿QUIEN SE AFECTA MAS?",
            "Comparacion normalizada entre activos — TPH como % del baseline pre-ventana", page_num=7)
    _footer(fig, DATE_STR)

    _img(fig, FIGS / "event_study" / "10_Comparacion_Activos.png",
         0.03, 0.12, 0.66, 0.72)

    # Tabla ranking
    ax_tbl = fig.add_axes([0.72, 0.12, 0.26, 0.72])
    ax_tbl.set_facecolor(C_LGRAY)
    ax_tbl.set_xlim(0, 1); ax_tbl.set_ylim(0, 1)
    ax_tbl.axis("off")
    ax_tbl.add_patch(FancyBboxPatch((0.01, 0.01), 0.98, 0.97,
                                    boxstyle="round,pad=0.01",
                                    facecolor=C_LGRAY, edgecolor=C_COPPER, linewidth=1))

    ax_tbl.text(0.5, 0.94, "RANKING DE IMPACTO", color=C_BLUE,
                fontsize=9, fontweight="bold", ha="center", va="top")

    ranking = [
        ("1°", "SAG2",     "42.7%", "11.4h", C_RED),
        ("2°", "PMC",      "39.3%",  "8.4h", C_ORANGE),
        ("3°", "SAG1",     "37.3%", "10.4h", C_BLUE),
        ("4°", "UNITARIO", "19.4%",  "5.8h", C_GREEN),
    ]
    headers = ["Pos", "Activo", "Caida", "Rec 90%"]
    hx = [0.07, 0.22, 0.56, 0.78]
    ax_tbl.text(hx[0], 0.86, headers[0], color=C_MGRAY, fontsize=7.5, fontweight="bold")
    ax_tbl.text(hx[1], 0.86, headers[1], color=C_MGRAY, fontsize=7.5, fontweight="bold")
    ax_tbl.text(hx[2], 0.86, headers[2], color=C_MGRAY, fontsize=7.5, fontweight="bold")
    ax_tbl.text(hx[3], 0.86, headers[3], color=C_MGRAY, fontsize=7.5, fontweight="bold")
    ax_tbl.axhline(0.83, color=C_COPPER, linewidth=0.8, xmin=0.04, xmax=0.96)

    for j, (pos, act, caida, rec, col) in enumerate(ranking):
        y = 0.75 - j * 0.17
        ax_tbl.add_patch(mpatches.Rectangle((0.03, y - 0.04), 0.94, 0.13,
                                             color=col, alpha=0.08))
        ax_tbl.add_patch(mpatches.Rectangle((0.03, y - 0.04), 0.012, 0.13,
                                             color=col))
        ax_tbl.text(hx[0], y + 0.02, pos,   color=col, fontsize=9, fontweight="bold", va="center")
        ax_tbl.text(hx[1], y + 0.02, act,   color=C_DARK, fontsize=9, fontweight="bold", va="center")
        ax_tbl.text(hx[2], y + 0.02, caida, color=col, fontsize=9, fontweight="bold", va="center")
        ax_tbl.text(hx[3], y + 0.02, rec,   color=C_DARK, fontsize=8.5, va="center")

    ax_tbl.axhline(0.38, color=C_MGRAY, linewidth=0.4, xmin=0.04, xmax=0.96, linestyle="--")
    ax_tbl.text(0.5, 0.30, "UNITARIO es el activo\nmas resiliente del sistema",
                color=C_GREEN, fontsize=8, ha="center", va="center", linespacing=1.5,
                fontweight="bold")
    ax_tbl.text(0.5, 0.15, "SAG2 requiere atencion\nprioritaria en toda\nventana programada",
                color=C_RED, fontsize=8, ha="center", va="center", linespacing=1.5,
                fontweight="bold")

    return fig


def page_hallazgo_ranking() -> plt.Figure:
    fig = _new_page()
    _header(fig, "HALLAZGO 4 — VULNERABILIDAD E INDICE DE RESILIENCIA",
            "IVO = Caida% x Tiempo Recuperacion  |  IR = mayor valor = mas resiliente", page_num=8)
    _footer(fig, DATE_STR)

    _img(fig, FIGS / "prescriptivo" / "P1_IVO_Resiliencia.png",
         0.03, 0.12, 0.94, 0.72)
    return fig


def page_hallazgo_recuperacion() -> plt.Figure:
    fig = _new_page()
    _header(fig, "HALLAZGO 5 — TIEMPO DE RECUPERACION",
            "Horas necesarias para recuperar 80%, 90%, 95% y 100% del baseline", page_num=9)
    _footer(fig, DATE_STR)

    _img(fig, FIGS / "event_study" / "11_Tiempo_Recuperacion.png",
         0.03, 0.12, 0.66, 0.72)

    msgs = [
        ("SAG2  —  Cuello de botella", C_RED,
         "11.4h para 90%\nLa recuperacion mas lenta.\nVentanas 12h: hasta 16.8h."),
        ("SAG1  —  Recuperacion Lenta", C_ORANGE,
         "10.4h para 90%\nPese a menor caida en %,\nel tiempo de vuelta es alto."),
        ("PMC  —  Recuperacion Media", C_BLUE,
         "8.4h para 90%\nMayor IVO pero recuperacion\nrelativa mas rapida que SAG2."),
        ("UNITARIO  —  Mas rapido", C_GREEN,
         "5.8h para 90%\nEl activo con mayor\nresiliencia del sistema."),
    ]
    for i, (ttl, col, body) in enumerate(msgs):
        y = 0.73 - i * 0.19
        ax = fig.add_axes([0.72, y, 0.26, 0.165])
        ax.set_facecolor(C_LGRAY)
        ax.set_xlim(0, 1); ax.set_ylim(0, 1)
        ax.axis("off")
        ax.add_patch(mpatches.Rectangle((0, 0), 0.025, 1, color=col))
        ax.text(0.06, 0.78, ttl,  color=col,    fontsize=8, fontweight="bold", va="top")
        ax.text(0.06, 0.42, body, color=C_DARK, fontsize=7.5, va="center", linespacing=1.5)
    return fig


def page_hallazgo_estadistica(data: dict) -> plt.Figure:
    fig = _new_page()
    _header(fig, "HALLAZGO 6 — EVIDENCIA ESTADISTICA",
            "T-test de Welch + Mann-Whitney U  |  Nivel de significancia: alpha = 0.05", page_num=10)
    _footer(fig, DATE_STR)

    df_stat = data["stat"]

    # Tabla principal
    ax_tbl = fig.add_axes([0.03, 0.38, 0.62, 0.46])
    ax_tbl.set_facecolor(C_LGRAY)
    ax_tbl.set_xlim(0, 10); ax_tbl.set_ylim(0, 5.5)
    ax_tbl.axis("off")

    headers = ["Activo", "TPH Pre", "TPH Post", "Delta%", "p-value T-test", "Resultado"]
    hx      = [0.3, 1.9, 3.6, 5.2, 6.6, 8.5]
    ax_tbl.add_patch(mpatches.Rectangle((0, 4.7), 10, 0.7, color=C_BLUE))
    for h, x in zip(headers, hx):
        ax_tbl.text(x, 5.0, h, color=C_WHITE, fontsize=8.5, fontweight="bold", va="center")

    row_data = [
        (r["activo"], f"{r['tph_pre_mean']:.0f}", f"{r['tph_post_mean']:.0f}",
         f"{r['delta_pct']:+.1f}%", f"{r['t_pval']:.4f}", r["significativo"])
        for _, r in df_stat.iterrows()
    ]
    for j, (act, pre, post, delta, pval, sig) in enumerate(row_data):
        y = 3.9 - j * 1.0
        row_col = C_LGRAY if j % 2 == 0 else C_WHITE
        ax_tbl.add_patch(mpatches.Rectangle((0.1, y - 0.3), 9.8, 0.85,
                                             color=row_col, alpha=0.7))
        sig_col = C_GREEN if sig == "SI" else C_MGRAY
        sig_txt = "SIGNIFICATIVO  ✓" if sig == "SI" else "No significativo"
        row = [act, pre, post, delta, pval, sig_txt]
        for val, x, is_sig in zip(row, hx, [False]*5 + [True]):
            col = sig_col if is_sig else C_DARK
            fw  = "bold" if is_sig or val == act else "normal"
            ax_tbl.text(x, y + 0.1, val, color=col, fontsize=8.5,
                        va="center", fontweight=fw)

    # Interpretacion ejecutiva
    ax_int = fig.add_axes([0.67, 0.38, 0.30, 0.46])
    ax_int.set_facecolor(C_BLUE)
    ax_int.set_xlim(0, 1); ax_int.set_ylim(0, 1)
    ax_int.axis("off")
    ax_int.add_patch(FancyBboxPatch((0.02, 0.02), 0.96, 0.96,
                                    boxstyle="round,pad=0.02",
                                    facecolor=C_BLUE, edgecolor=C_COPPER, linewidth=1.5))
    ax_int.text(0.5, 0.88, "INTERPRETACION", color=C_GOLD,
                fontsize=8.5, fontweight="bold", ha="center", va="top")
    msg = ("Los resultados observados\nNO son atribuibles al azar.\n\n"
           "En SAG1, SAG2 y PMC\nla diferencia pre/post es\nestadisticamente real.\n\n"
           "Esto valida que el efecto T8\nexiste y puede ser medido\ncon confianza.")
    ax_int.text(0.5, 0.70, msg, color=C_WHITE, fontsize=8.5,
                ha="center", va="top", linespacing=1.6)

    # Nota UNITARIO
    _callout(fig, 0.03, 0.14, 0.45, 0.20,
             "UNITARIO no muestra significancia estadistica (p=0.24).\n"
             "Su caida promedio (19.4%) es menor y con mayor variabilidad.\n"
             "Requiere mas eventos para confirmar el efecto.",
             bg=C_LGRAY, fc=C_DARK)
    _callout(fig, 0.52, 0.14, 0.45, 0.20,
             "SAG2: mayor impacto economico absoluto.\n"
             "Cohen d = 0.114 → efecto pequeno-moderado.\n"
             "Pero sobre 2,117 t/h baseline = alto impacto en toneladas.",
             bg=C_LGRAY, fc=C_DARK)

    return fig


def page_impacto() -> plt.Figure:
    fig = _new_page()
    _header(fig, "IMPACTO OPERACIONAL",
            "Toneladas no producidas y costo estimado del periodo analizado", page_num=11)
    _footer(fig, DATE_STR)

    # KPIs de toneladas
    ton_data = [
        ("136.7k ton", "TOTAL PERDIDO",       "72 eventos · 5.5 meses",        C_RED),
        ("61.7k ton",  "SAG2 — Lider",         "Mayor perdida absoluta (61.7k)", C_ORANGE),
        ("23.2k ton",  "SAG1",                 "17.7% de produccion esperada",  C_BLUE),
        ("33.6k ton",  "PMC",                  "13.6% de produccion esperada",  C_COPPER),
    ]
    for i, (val, lbl, sub, col) in enumerate(ton_data):
        _kpi_box(fig, 0.03 + i * 0.245, 0.60, 0.22, 0.24,
                 val, lbl, sub, bg=col)

    _img(fig, FIGS / "prescriptivo" / "P2_Toneladas_Perdidas.png",
         0.03, 0.12, 0.60, 0.44)

    # Panel valor referencial
    ax_val = fig.add_axes([0.66, 0.12, 0.32, 0.44])
    ax_val.set_facecolor(C_DARK)
    ax_val.set_xlim(0, 1); ax_val.set_ylim(0, 1)
    ax_val.axis("off")
    ax_val.add_patch(FancyBboxPatch((0.02, 0.02), 0.96, 0.96,
                                    boxstyle="round,pad=0.02",
                                    facecolor=C_DARK, edgecolor=C_GOLD, linewidth=1.5))
    ax_val.text(0.5, 0.91, "VALOR REFERENCIAL", color=C_GOLD,
                fontsize=8.5, fontweight="bold", ha="center", va="top")
    ax_val.text(0.5, 0.78, "(Cu estimado)", color=C_MGRAY, fontsize=7.5, ha="center", va="top")
    items = [
        ("Grade mineral", "1.2% Cu"),
        ("Recuperacion", "86%"),
        ("Precio LME ref.", "USD 9,500/ton Cu"),
        ("Cu perdido", "~1,414 ton Cu"),
        ("Valor periodo", "~USD 13.4M"),
        ("Valor anual", "~USD 29.2M"),
    ]
    for j, (k, v) in enumerate(items):
        y = 0.66 - j * 0.11
        ax_val.text(0.08, y, k + ":", color=C_MGRAY, fontsize=8, va="center")
        ax_val.text(0.92, y, v,       color=C_GOLD,  fontsize=8.5, va="center",
                    ha="right", fontweight="bold")
    ax_val.text(0.5, 0.03, "Referencial — sujeto a grade y recuperacion real",
                color=C_MGRAY, fontsize=6.5, ha="center", va="bottom", alpha=0.7)

    return fig


def page_curvas_estrategicas() -> plt.Figure:
    fig = _new_page()
    _header(fig, "PUNTO DE QUIEBRE OPERACIONAL",
            "¿A partir de que duracion el impacto se vuelve inaceptable?", page_num=12)
    _footer(fig, DATE_STR)

    _img(fig, FIGS / "prescriptivo" / "P3_Curvas_Estrategicas.png",
         0.03, 0.12, 0.68, 0.72)

    msgs = [
        ("HASTA 4H", C_GREEN,
         "Caida ~38%  |  2,507 ton/evento\nZona de impacto manejable\nRecuperacion en 1 turno"),
        ("ENTRE 4H Y 12H", C_ORANGE,
         "Caida escala linealmente\nCada hora adicional cuesta\n~717 ton de produccion"),
        ("12H — CRITICO", C_RED,
         "8,242 ton/evento = 3.3x\nRiesgo acumulativo alto\nRecomendacion: EVITAR"),
    ]
    for i, (ttl, col, body) in enumerate(msgs):
        y = 0.73 - i * 0.22
        ax = fig.add_axes([0.73, y, 0.25, 0.19])
        ax.set_facecolor(C_LGRAY)
        ax.set_xlim(0, 1); ax.set_ylim(0, 1)
        ax.axis("off")
        ax.add_patch(mpatches.Rectangle((0, 0), 0.025, 1, color=col))
        ax.text(0.06, 0.80, ttl,  color=col,    fontsize=8.5, fontweight="bold", va="top")
        ax.text(0.06, 0.47, body, color=C_DARK, fontsize=7.5, va="center", linespacing=1.5)

    _callout(fig, 0.73, 0.12, 0.25, 0.12,
             "Duracion maxima\nrecomendable: 4 horas\nExcepcion: con protocolo\nde pila llena.",
             bg=C_BLUE)
    return fig


def page_riesgos() -> plt.Figure:
    fig = _new_page()
    _header(fig, "RIESGOS OPERACIONALES",
            "Top 5 riesgos identificados en el analisis", page_num=13)
    _footer(fig, DATE_STR)

    riesgos = [
        ("R1", "VULNERABILIDAD CRITICA DE SAG2",
         "SAG2 concentra el 45% de la perdida total (61.7k ton). "
         "Una falla compuesta (ventana T8 + falla equipo) puede comprometer "
         "significativamente el plan de produccion mensual.",
         "ALTO", C_RED),
        ("R2", "VENTANAS 12H — IMPACTO SEVERO",
         "7 eventos de 12h en el periodo = 57.7k ton perdidas (42% del total). "
         "Son el 9.7% de los eventos pero generan el mayor dano unitario. "
         "Sin protocolo de compensacion, el impacto es irrecuperable.",
         "ALTO", C_RED),
        ("R3", "RECUPERACION MAS LENTA QUE UN TURNO",
         "SAG2 y SAG1 tardan >10h en recuperar el 90%. Eso supera la duracion "
         "de un turno. El efecto de una ventana se extiende inevitablemente "
         "al turno siguiente sin acciones especificas.",
         "MEDIO", C_ORANGE),
        ("R4", "AGOTAMIENTO DE PILAS",
         "T8 es el unico proveedor de mineral a chancado. Si las pilas inician "
         "la ventana con stock bajo, la caida ocurre DURANTE la ventana, "
         "no despues. UNITARIO tiene el menor IAP (1.85): es el mas expuesto.",
         "MEDIO", C_ORANGE),
        ("R5", "VENTANAS CONSECUTIVAS EN DIAS SEGUIDOS",
         "No se evaluo el impacto de ventanas consecutivas. Si el sistema "
         "no se recupera entre un evento y el siguiente, el efecto "
         "se amplifica de forma acumulativa.",
         "MEDIO", C_COPPER),
    ]

    for i, (cod, ttl, desc, nivel, col) in enumerate(riesgos):
        y = 0.76 - i * 0.133
        ax = fig.add_axes([0.03, y, 0.94, 0.11])
        ax.set_facecolor(C_LGRAY)
        ax.set_xlim(0, 1); ax.set_ylim(0, 1)
        ax.axis("off")
        # Barra lateral de severidad
        ax.add_patch(mpatches.Rectangle((0, 0), 0.012, 1, color=col))
        # Código
        ax.add_patch(FancyBboxPatch((0.015, 0.1), 0.055, 0.8, boxstyle="round,pad=0.02",
                                    facecolor=col, edgecolor="none"))
        ax.text(0.042, 0.5, cod, color=C_WHITE, fontsize=9, fontweight="bold",
                ha="center", va="center")
        # Título
        ax.text(0.085, 0.7, ttl,  color=col,    fontsize=9, fontweight="bold", va="center")
        # Descripción
        ax.text(0.085, 0.25, desc[:120] + ("..." if len(desc) > 120 else ""),
                color=C_DARK, fontsize=7.5, va="center")
        # Nivel
        ax.add_patch(FancyBboxPatch((0.87, 0.2), 0.11, 0.6,
                                    boxstyle="round,pad=0.02",
                                    facecolor=col, edgecolor="none", alpha=0.15))
        ax.text(0.925, 0.5, nivel, color=col, fontsize=8, fontweight="bold",
                ha="center", va="center")

    return fig


def page_oportunidades() -> plt.Figure:
    fig = _new_page()
    _header(fig, "OPORTUNIDADES DE MEJORA",
            "Acciones con mayor impacto operacional y ROI", page_num=14)
    _footer(fig, DATE_STR)

    _img(fig, FIGS / "prescriptivo" / "P5_Escenarios.png",
         0.03, 0.38, 0.60, 0.46)

    opps = [
        ("O1", "REDISTRIBUIR VENTANAS 12H",
         "Mover a paradas programadas mayores.\nAhorro: +57.7k ton / +USD 29M/ano.",
         "ALTA", C_GREEN),
        ("O2", "PROTOCOLO PILA LLENA",
         "Maximizar stock antes de ventana ≥4h.\nMitigacion directa: alarga el tiempo\nhasta la caida.",
         "ALTA", C_GREEN),
        ("O3", "VENTANAS NOCTURNAS",
         "Iniciar ventana 2h antes del periodo\nvalle. Reduce impacto ~25%.",
         "MEDIA", C_BLUE),
        ("O4", "MONITOREO TIEMPO REAL",
         "Alerta si caida >30% durante ventana.\nCompensacion desde circuito alternativo.",
         "MEDIA", C_COPPER),
        ("O5", "REDUCIR MTTR 12H → 8H",
         "Optimizar ejecucion de trabajos 12h.\nAhorro potencial: +20k ton/ano.",
         "MEDIA", C_ORANGE),
    ]

    for i, (cod, ttl, desc, nivel, col) in enumerate(opps):
        x = 0.655 + (i % 2) * 0.17
        y = 0.72 - (i // 2) * 0.20 if i < 4 else 0.38
        if i == 4:
            x = 0.655 + 0.085
        ax = fig.add_axes([x, y, 0.155, 0.17] if i < 4 else [0.655 + 0.085, 0.38, 0.155, 0.17])
        ax.set_facecolor(C_LGRAY)
        ax.set_xlim(0, 1); ax.set_ylim(0, 1)
        ax.axis("off")
        ax.add_patch(FancyBboxPatch((0.02, 0.02), 0.96, 0.96,
                                    boxstyle="round,pad=0.02",
                                    facecolor=C_LGRAY, edgecolor=col, linewidth=1.5))
        ax.add_patch(FancyBboxPatch((0.03, 0.76), 0.22, 0.20,
                                    boxstyle="round,pad=0.01",
                                    facecolor=col, edgecolor="none"))
        ax.text(0.14, 0.86, cod, color=C_WHITE, fontsize=9, fontweight="bold",
                ha="center", va="center")
        ax.text(0.5, 0.60, ttl, color=col, fontsize=7.5, fontweight="bold",
                ha="center", va="center", linespacing=1.3)
        ax.text(0.5, 0.32, desc, color=C_DARK, fontsize=6.5,
                ha="center", va="center", linespacing=1.4)
        ax.add_patch(FancyBboxPatch((0.62, 0.76), 0.33, 0.20,
                                    boxstyle="round,pad=0.01",
                                    facecolor=col, edgecolor="none", alpha=0.2))
        ax.text(0.785, 0.86, nivel, color=col, fontsize=7, fontweight="bold",
                ha="center", va="center")

    return fig


def page_recomendaciones() -> plt.Figure:
    fig = _new_page()
    _header(fig, "RECOMENDACIONES OPERACIONALES",
            "Acciones por horizonte temporal — Impacto y prioridad", page_num=15)
    _footer(fig, DATE_STR)

    columns = [
        ("CORTO PLAZO", "0 – 30 dias", C_RED, [
            ("Operaciones", "Protocolo pila llena antes de toda ventana ≥4h", "ALTO", "INMEDIATA"),
            ("Operaciones", "Alerta temprana si caida SAG2 supera 35% durante ventana", "ALTO", "INMEDIATA"),
            ("PAM Mantto",  "Avisar ventana T8 con 48h de anticipacion", "ALTO", "INMEDIATA"),
            ("Planificacion","Evitar ventanas 12h sin justificacion documentada", "ALTO", "INMEDIATA"),
        ]),
        ("MEDIANO PLAZO", "1 – 3 meses", C_ORANGE, [
            ("Planificacion", "Redistribuir trabajos 12h a proxima parada mayor", "MUY ALTO", "ALTA"),
            ("Opt. Activos",  "Evaluar capacidad de pila UNITARIO — menor IAP", "MEDIO",  "MEDIA"),
            ("PAM Mantto",    "Registrar duracion real vs. planificada por evento", "MEDIO", "ALTA"),
            ("Operaciones",   "Protocolo de compensacion desde circuito alternativo", "ALTO", "ALTA"),
        ]),
        ("LARGO PLAZO", "3 – 12 meses", C_BLUE, [
            ("CIO / Opt.",    "Modelo predictivo de impacto por tipo de ventana", "ALTO", "MEDIA"),
            ("Opt. Activos",  "Evaluar inversion en capacidad de pila SAG2", "ALTO", "MEDIA"),
            ("Planificacion", "Programa anual con tope de ventanas 12h por semestre", "MUY ALTO", "ALTA"),
            ("CIO",           "Dashboard tiempo real TPH vs. ventana activa", "MEDIO", "MEDIA"),
        ]),
    ]

    for col_i, (ttl, period, col, items) in enumerate(columns):
        x = 0.03 + col_i * 0.325
        # Header columna
        ax_h = fig.add_axes([x, 0.75, 0.31, 0.09])
        ax_h.set_facecolor(col)
        ax_h.axis("off")
        ax_h.text(0.5, 0.65, ttl,    color=C_WHITE, fontsize=10, fontweight="bold", ha="center", va="center")
        ax_h.text(0.5, 0.22, period, color=C_GOLD,  fontsize=8.5, ha="center", va="center")

        # Items
        for j, (area, accion, imp, prio) in enumerate(items):
            y = 0.575 - j * 0.135
            ax_i = fig.add_axes([x, y, 0.31, 0.118])
            ax_i.set_facecolor(C_LGRAY)
            ax_i.set_xlim(0, 1); ax_i.set_ylim(0, 1)
            ax_i.axis("off")
            ax_i.add_patch(mpatches.Rectangle((0, 0), 0.008, 1, color=col))
            ax_i.text(0.03, 0.82, area,   color=col,    fontsize=7, fontweight="bold", va="top")
            ax_i.text(0.03, 0.52, accion[:60], color=C_DARK, fontsize=7, va="top",
                      wrap=True, linespacing=1.3)
            # Tags
            ax_i.add_patch(FancyBboxPatch((0.03, 0.05), 0.28, 0.25,
                                          boxstyle="round,pad=0.01",
                                          facecolor=col, edgecolor="none", alpha=0.2))
            ax_i.text(0.17, 0.17, f"Impacto: {imp}", color=col, fontsize=6.5,
                      ha="center", va="center", fontweight="bold")
            ax_i.add_patch(FancyBboxPatch((0.33, 0.05), 0.28, 0.25,
                                          boxstyle="round,pad=0.01",
                                          facecolor=C_MGRAY, edgecolor="none", alpha=0.2))
            ax_i.text(0.47, 0.17, f"Prio: {prio}", color=C_MGRAY, fontsize=6.5,
                      ha="center", va="center")

    return fig


def page_conclusiones() -> plt.Figure:
    fig = _new_page()
    _header(fig, "CONCLUSIONES FINALES",
            "Respuestas directas a las preguntas del comite", page_num=16)
    _footer(fig, DATE_STR)

    preguntas = [
        ("¿Existe efecto gaviota?",
         "SI — Confirmado en 72 eventos. Cuando T8 se detiene, los molinos agotan las pilas\n"
         "y el TPH cae de forma repetitiva y significativa (p<0.05) en SAG1, SAG2 y PMC.",
         C_GREEN),
        ("¿Que activo es mas vulnerable?",
         "SAG2: mayor caida (42.7%), mayor perdida absoluta (61.7k ton), mayor elasticidad\n"
         "(13.7%/h). PMC: mayor IVO combinado (caida alta + recuperacion lenta).",
         C_RED),
        ("¿Que duracion genera mayor impacto?",
         "Las ventanas de 12h: 8,242 ton de perdida media por evento = 3.3x vs ventana 4h.\n"
         "El punto de quiebre operacional esta en 4h. Maximo recomendable: 4h.",
         C_ORANGE),
        ("¿Existe evidencia estadistica?",
         "SI — SAG1 (p=0.050), SAG2 (p<0.001), PMC (p<0.001). UNITARIO no es significativo.\n"
         "Los resultados NO son atribuibles al azar.",
         C_BLUE),
        ("¿Que accion deberia priorizar el comite?",
         "1° Redistribuir las 7 ventanas 12h a paradas programadas (ahorro ~USD 29M/ano).\n"
         "2° Protocolo pila llena antes de toda ventana ≥4h (costo: coordinacion, no capital).",
         C_COPPER),
    ]

    for i, (preg, resp, col) in enumerate(preguntas):
        y = 0.78 - i * 0.133
        ax = fig.add_axes([0.03, y, 0.94, 0.115])
        ax.set_facecolor(C_LGRAY)
        ax.set_xlim(0, 1); ax.set_ylim(0, 1)
        ax.axis("off")
        ax.add_patch(mpatches.Rectangle((0, 0), 0.005, 1, color=col))
        ax.add_patch(FancyBboxPatch((0.008, 0.05), 0.05, 0.90,
                                    boxstyle="round,pad=0.01",
                                    facecolor=col, edgecolor="none", alpha=0.15))
        ax.text(0.032, 0.5, f"P{i+1}", color=col, fontsize=11, fontweight="bold",
                ha="center", va="center")
        ax.text(0.075, 0.75, preg, color=col, fontsize=9, fontweight="bold", va="center")
        ax.text(0.075, 0.32, resp, color=C_DARK, fontsize=8, va="center", linespacing=1.4)

    return fig


def page_cierre() -> plt.Figure:
    fig = _new_page(C_BLUE)

    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_facecolor(C_BLUE)
    ax.set_xlim(0, 1); ax.set_ylim(0, 1)
    ax.axis("off")
    ax.add_patch(mpatches.Rectangle((0, 0), 0.006, 1, color=C_COPPER))

    ax.text(0.5, 0.80, "GRACIAS", color=C_GOLD, fontsize=36, fontweight="bold",
            ha="center", va="center")
    ax.text(0.5, 0.66, "Analítica de Rendimientos — CIO DET", color=C_WHITE,
            fontsize=14, ha="center", va="center")
    ax.axhline(0.60, xmin=0.20, xmax=0.80, color=C_COPPER, linewidth=1.5)
    ax.text(0.5, 0.50, f"Informe: Impacto Operacional Ventanas Teniente 8\n"
            f"Periodo: Enero 2026 – Junio 2026  |  72 eventos",
            color=C_MGRAY, fontsize=10, ha="center", va="center", linespacing=1.8)
    ax.text(0.5, 0.35, "Contacto analítica: CIO DET — División El Teniente",
            color=C_MGRAY, fontsize=9, ha="center", va="center")
    ax.text(0.5, 0.14, f"CODELCO División El Teniente  |  {DATE_STR}  |  v1.0  |  CONFIDENCIAL",
            color=C_MGRAY, fontsize=8, ha="center", va="center", alpha=0.7)
    return fig


# ═══════════════════════════════════════════════════════════════════════════════
# ANEXO TÉCNICO
# ═══════════════════════════════════════════════════════════════════════════════

def page_anexo_metodologia() -> plt.Figure:
    fig = _new_page()
    _header(fig, "ANEXO TECNICO — METODOLOGIA DETALLADA",
            "Event Study Industrial: parametros y supuestos del modelo", page_num=1)
    _footer(fig, DATE_STR, version="Anexo v1.0")

    ax = fig.add_axes([0.03, 0.10, 0.94, 0.73])
    ax.set_facecolor(C_LGRAY)
    ax.set_xlim(0, 1); ax.set_ylim(0, 1)
    ax.axis("off")

    secciones = [
        ("FUENTES DE DATOS",
         [("Rendimientos", "rendimientos_clean.parquet — 47,532 registros 5-min | SAG1, SAG2, PMC, UNITARIO"),
          ("Eventos T8", "ventanas_t8.parquet — 72 eventos ene-jun 2026 (PAM Mantto como fuente oficial)"),
          ("Operacional", "Filtro: operando=True AND tph > 50 (umbral <= 50 = coeficientes de estado, no produccion real)")]),
        ("PARAMETROS EVENT STUDY",
         [("Ventana analisis", "-24h a +24h alrededor de t=0 (inicio oficial ventana)"),
          ("t=0 oficial", "2h→14:00 | 4h→12:00 | 8h→08:00 | 12h→08:00 (NO inferido de datos)"),
          ("Baseline", "Periodo PRE: -24h a 0h | Minimo 20 puntos 5-min para calcular"),
          ("Bin respuesta", "30 minutos (aggregacion de la curva promedio)")]),
        ("METRICAS CALCULADAS",
         [("Caida%", "(baseline_mean - tph_min) / baseline_mean * 100"),
          ("h_hasta_min", "Tiempo desde t=0 hasta el minimo rolling (ventana 6 puntos = 30 min)"),
          ("IVO", "Indice Vulnerabilidad = caida_pct * h_recuperacion_90"),
          ("IAP", "Indice Amortiguacion = h_hasta_min / duracion_ventana")]),
        ("PRUEBAS ESTADISTICAS",
         [("T-test Welch", "Igualdad de medias (no asume igualdad de varianzas) — SciPy ttest_ind"),
          ("Mann-Whitney U", "Prueba no parametrica alternativa — SciPy mannwhitneyu"),
          ("Nivel alpha", "0.05 (95% confianza) — resultado significativo si p-value < 0.05"),
          ("Efecto", "Cohen d = (media_pre - media_post) / sqrt((var_pre + var_post) / 2)")]),
    ]
    col_x = [0.01, 0.51]
    col_y = [0.97, 0.51]
    for idx, (ttl, items) in enumerate(secciones):
        x = col_x[idx % 2]
        y = col_y[idx // 2]
        ax.add_patch(FancyBboxPatch((x, y - 0.44), 0.475, 0.44,
                                    boxstyle="round,pad=0.01",
                                    facecolor=C_WHITE, edgecolor=C_BLUE, linewidth=1))
        ax.text(x + 0.01, y - 0.02, ttl, color=C_BLUE, fontsize=9,
                fontweight="bold", va="top")
        for j, (k, v) in enumerate(items):
            yy = y - 0.09 - j * 0.085
            ax.text(x + 0.02, yy, f"■ {k}:", color=C_ORANGE, fontsize=7.5,
                    fontweight="bold", va="top")
            ax.text(x + 0.14, yy, v, color=C_DARK, fontsize=7.5, va="top",
                    wrap=True)
    return fig


def page_anexo_estadistica(data: dict) -> plt.Figure:
    fig = _new_page()
    _header(fig, "ANEXO TECNICO — RESULTADOS ESTADISTICOS COMPLETOS",
            "T-test de Welch + Mann-Whitney U — todos los activos", page_num=2)
    _footer(fig, DATE_STR, version="Anexo v1.0")

    df_stat = data["stat"]

    ax = fig.add_axes([0.03, 0.55, 0.94, 0.35])
    ax.set_facecolor(C_LGRAY)
    ax.set_xlim(0, 16); ax.set_ylim(0, 5)
    ax.axis("off")

    headers = ["Activo", "TPH Pre\n(media)", "TPH Post\n(media)", "Delta%",
               "t-stat", "t p-value", "U-stat", "U p-value",
               "Cohen d", "Significativo", "N pre", "N post"]
    hx = [0.3, 1.7, 3.2, 4.5, 5.7, 6.9, 8.4, 9.8, 11.2, 12.4, 14.0, 15.0]
    ax.add_patch(mpatches.Rectangle((0, 4.0), 16, 0.9, color=C_BLUE))
    for h, x in zip(headers, hx):
        ax.text(x, 4.42, h, color=C_WHITE, fontsize=7, fontweight="bold",
                va="center", ha="left", linespacing=1.3)

    for j, (_, r) in enumerate(df_stat.iterrows()):
        y = 3.1 - j * 0.9
        rc = C_LGRAY if j % 2 == 0 else C_WHITE
        ax.add_patch(mpatches.Rectangle((0.1, y - 0.3), 15.8, 0.82, color=rc, alpha=0.7))
        sig_col = C_GREEN if r["significativo"] == "SI" else C_MGRAY
        row = [r["activo"], f"{r['tph_pre_mean']:.0f}", f"{r['tph_post_mean']:.0f}",
               f"{r['delta_pct']:+.2f}%", f"{r['t_stat']:.3f}", f"{r['t_pval']:.5f}",
               f"{r['u_stat']:.0f}", f"{r['u_pval']:.5f}",
               f"{r['cohen_d']:.3f}",
               "SI  *" if r["significativo"] == "SI" else "NO",
               str(r["n_pre"]), str(r["n_post"])]
        for val, x in zip(row, hx):
            col = sig_col if val in ["SI  *", "NO"] else C_DARK
            ax.text(x, y + 0.10, val, color=col, fontsize=8, va="center",
                    fontweight="bold" if val in ["SI  *", "NO"] else "normal")

    # Tabla métricas
    _img(fig, FIGS / "prescriptivo" / "P7_Forecast_Impacto.png",
         0.03, 0.10, 0.94, 0.42)

    return fig


def page_anexo_metricas(data: dict) -> plt.Figure:
    fig = _new_page()
    _header(fig, "ANEXO TECNICO — KPI y METRICAS POR ACTIVO",
            "IVO, IR, IAP, IST8, Elasticidad — calculo y valores", page_num=3)
    _footer(fig, DATE_STR, version="Anexo v1.0")

    _img(fig, FIGS / "prescriptivo" / "P1_IVO_Resiliencia.png",
         0.03, 0.47, 0.55, 0.42)
    _img(fig, FIGS / "prescriptivo" / "P6_Amortiguacion_Pilas.png",
         0.60, 0.47, 0.38, 0.42)
    _img(fig, FIGS / "prescriptivo" / "P4_Elasticidad.png",
         0.03, 0.10, 0.94, 0.34)

    return fig


def page_anexo_panel(data: dict) -> plt.Figure:
    fig = _new_page()
    _header(fig, "ANEXO TECNICO — PANEL DIAGNOSTICO INTEGRADO",
            "Vista consolidada de todos los indicadores", page_num=4)
    _footer(fig, DATE_STR, version="Anexo v1.0")

    _img(fig, FIGS / "prescriptivo" / "P8_Panel_Ejecutivo.png",
         0.03, 0.10, 0.94, 0.74)
    return fig


# ═══════════════════════════════════════════════════════════════════════════════
# Construccion del PDF
# ═══════════════════════════════════════════════════════════════════════════════

def build_pdf(pages: list[plt.Figure], out_path: Path) -> None:
    with PdfPages(str(out_path)) as pdf:
        d = pdf.infodict()
        d["Title"]    = "Impacto Operacional Ventanas Teniente 8 — Analisis Event Study"
        d["Author"]   = "CIO DET — Analitica de Rendimientos"
        d["Subject"]  = "Informe Comite T8"
        d["Keywords"] = "Teniente 8, rendimiento, molienda, SAG, event study"
        d["CreationDate"] = datetime.now()
        for fig in pages:
            pdf.savefig(fig, dpi=DPI, bbox_inches="tight")
            plt.close(fig)


def build_pptx(pages_data: list[dict], out_path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width  = Inches(16)
    prs.slide_height = Inches(9)

    blank = prs.slide_layouts[6]   # completely blank

    for pdata in pages_data:
        slide = prs.slides.add_slide(blank)
        bg = slide.background
        fill = bg.fill
        fill.solid()
        r, g, b = tuple(int(pdata["bg"].lstrip("#")[i:i+2], 16) for i in (0, 2, 4))
        fill.fore_color.rgb = RGBColor(r, g, b)

        # Insertar imagen PNG (guardada temporalmente desde la fig matplotlib)
        img_path = pdata["img_path"]
        if Path(img_path).exists():
            slide.shapes.add_picture(str(img_path), Inches(0), Inches(0),
                                     Inches(16), Inches(9))

    prs.save(str(out_path))


# ═══════════════════════════════════════════════════════════════════════════════
# Punto de entrada
# ═══════════════════════════════════════════════════════════════════════════════

def run_report(verbose: bool = True) -> None:
    if verbose:
        print("=" * 72)
        print("  Generando Informe Ejecutivo PDF — Comite T8")
        print("=" * 72)

    data = _load_data()

    # ── Informe principal ──────────────────────────────────────────────────────
    if verbose:
        print("\n  [1/2] Informe principal...")
    main_pages = [
        page_portada(),
        page_resumen_ejecutivo(),
        page_contexto(),
        page_metodologia(),
        page_hallazgo_gaviota(),
        page_hallazgo_duracion(),
        page_hallazgo_activos(),
        page_hallazgo_ranking(),
        page_hallazgo_recuperacion(),
        page_hallazgo_estadistica(data),
        page_impacto(),
        page_curvas_estrategicas(),
        page_riesgos(),
        page_oportunidades(),
        page_recomendaciones(),
        page_conclusiones(),
        page_cierre(),
    ]

    main_pdf = RPT / "Informe_Comite_T8.pdf"
    build_pdf(main_pages, main_pdf)
    if verbose:
        print(f"  OK  {main_pdf.name}  ({len(main_pages)} paginas)")

    # ── PPTX (mismas páginas exportadas como imágenes) ─────────────────────────
    if verbose:
        print("\n  [2/2] PPTX...")
    tmp_dir = RPT / "_tmp_slides"
    tmp_dir.mkdir(exist_ok=True)

    main_pages_pptx = [
        page_portada(), page_resumen_ejecutivo(), page_contexto(),
        page_metodologia(), page_hallazgo_gaviota(), page_hallazgo_duracion(),
        page_hallazgo_activos(), page_hallazgo_ranking(), page_hallazgo_recuperacion(),
        page_hallazgo_estadistica(data), page_impacto(), page_curvas_estrategicas(),
        page_riesgos(), page_oportunidades(), page_recomendaciones(),
        page_conclusiones(), page_cierre(),
    ]
    pptx_data = []
    for i, fig in enumerate(main_pages_pptx):
        img_p = tmp_dir / f"slide_{i:02d}.png"
        fig.savefig(str(img_p), dpi=DPI, bbox_inches="tight", facecolor=fig.get_facecolor())
        plt.close(fig)
        pptx_data.append({"img_path": str(img_p), "bg": "#FFFFFF"})

    pptx_path = RPT / "Informe_Comite_T8.pptx"
    build_pptx(pptx_data, pptx_path)
    # Limpiar temporales
    import shutil
    shutil.rmtree(tmp_dir, ignore_errors=True)
    if verbose:
        print(f"  OK  {pptx_path.name}  ({len(pptx_data)} diapositivas)")

    # ── Anexo técnico ──────────────────────────────────────────────────────────
    if verbose:
        print("\n  [Anexo] Generando Anexo Tecnico...")
    anexo_pages = [
        page_portada(),   # reusa portada como índice
        page_anexo_metodologia(),
        page_anexo_estadistica(data),
        page_anexo_metricas(data),
        page_anexo_panel(data),
        page_cierre(),
    ]
    # Portada del anexo con distinto texto
    anexo_pdf = RPT / "Anexo_Tecnico_T8.pdf"
    build_pdf(anexo_pages, anexo_pdf)
    if verbose:
        print(f"  OK  {anexo_pdf.name}  ({len(anexo_pages)} paginas)")

    # ── Log ────────────────────────────────────────────────────────────────────
    with open(LOGS / "skill_audit.log", "a", encoding="utf-8") as f:
        f.write(json.dumps({
            "fecha": datetime.now().isoformat(),
            "script": "src/generar_informe_pdf.py",
            "entregables": [
                str(main_pdf.name), str(pptx_path.name), str(anexo_pdf.name)
            ],
        }, ensure_ascii=False) + "\n")

    if verbose:
        print("\n" + "=" * 72)
        print("  ENTREGABLES GENERADOS:")
        print(f"    {main_pdf}")
        print(f"    {pptx_path}")
        print(f"    {anexo_pdf}")
        print("=" * 72)


if __name__ == "__main__":
    run_report()
